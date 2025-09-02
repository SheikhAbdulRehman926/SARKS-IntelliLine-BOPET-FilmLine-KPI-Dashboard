import os
import re
import base64
from datetime import datetime
from typing import Optional, Dict, Tuple, List

import numpy as np
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import streamlit as st
from io import BytesIO
from PIL import Image   # 👈 add this import with the others above

# ───────────────────────────────────────────────────────────────────────────────
# PAGE CONFIG
# ───────────────────────────────────────────────────────────────────────────────
tab_logo = Image.open("ChatGPT Image Sep 2, 2025, 12_39_59 PM.png")  # make sure this file is in the same folder as App_updated_SARKS.py

st.set_page_config(
    page_title="SARKS IntelliLine — BOPET FilmLine KPI Dashboard",
    page_icon=tab_logo,   # 👈 your uploaded logo as favicon
    layout="wide",
    initial_sidebar_state="expanded",
)


# =========================
# Brand & Colors
# =========================
BRAND      = "#F56023"
SLATE_970  = "#07090C"
SLATE_950  = "#0A0C10"
SLATE_900  = "#12161D"
SLATE_860  = "#151A22"
SLATE_850  = "#161B22"
SLATE_830  = "#171E29"
SLATE_820  = "#18202A"
SLATE_800  = "#1B2028"
SLATE_760  = "#222A36"
SLATE_720  = "#252C36"
TEXT       = "#EAEFF5"
MUTED      = "#96A0AD"
TEAL       = "#14B8A6"

LOGO_HEIGHT_PX = 140
HEADER_LOGO_COL_PX = max(260, min(520, int(LOGO_HEIGHT_PX * 10)))
HEADER_MIN_H = max(100, LOGO_HEIGHT_PX + 10)

COLORWAY = [
    "#60A5FA", "#34D399", "#F59E0B", "#A78BFA", "#F97316", "#EF4444",
    "#06B6D4", "#22C55E", "#D946EF", "#0EA5E9", "#F43F5E", "#84CC16",
]

# --- Developer & Footer constants ---
DEV_NAME  = "SARKS-Sheikh Abdul Rehman Bin Khalid Sharif"
DEV_ROLE  = "Developer / Artificial Intelligence Engineer • Internship at SARKS"
DEV_LINK  = "https://www.linkedin.com/in/sheikh-abdul-rehman-a14a982ab"

FOOTER_LEFT        = "© 2025 • SARKS IntelliLine — BOPET FilmLine KPI Dashboard  • Version-01 "
FOOTER_RIGHT_TEXT  = "Powered by SARKS"

# keep this for backward-compat with PDF/export code that still uses it
FOOTER_RIGHT_BRAND = FOOTER_RIGHT_TEXT
# (optional, if your PDF code embeds a link)
FOOTER_RIGHT_URL   = DEV_LINK

# =========================
# Utility helpers
# =========================
def _logo_b64(paths: Tuple[str, ...] = (
    "sarks-01.png", "SARKS-01.png", "./sarks-01.png", "./SARKS-01.png", "/mnt/data/sarks-01.png", "/mnt/data/SARKS-01.png"
)) -> Optional[str]:
    for p in paths:
        if os.path.exists(p):
            with open(p, "rb") as f:
                return base64.b64encode(f.read()).decode("utf-8")
    return None

def _logo_bytes() -> Optional[bytes]:
    for p in ("sarks-01.png", "SARKS-01.png", "./sarks-01.png", "./SARKS-01.png", "/mnt/data/sarks-01.png", "/mnt/data/SARKS-01.png"):
        if os.path.exists(p):
            with open(p, "rb") as f:
                return f.read()
    return None

def _duration_to_hours(x: str) -> Optional[float]:
    if not isinstance(x, str):
        return None
    s = x.strip()
    if re.fullmatch(r"\d{1,3}:\d{2}", s):
        h, m = s.split(":")
        try:
            return int(h) + int(m) / 60.0
        except Exception:
            return None
    return None

def _numeric(s: pd.Series, treat_as_time: bool = False) -> pd.Series:
    if treat_as_time:
        conv = s.apply(lambda v: _duration_to_hours(str(v)) if pd.notna(v) else np.nan)
        if conv.notna().sum() > 0:
            return pd.to_numeric(conv, errors="coerce")
    s2 = (
        s.astype(str)
         .str.replace(",", "", regex=False)
         .str.replace("%", "", regex=False)
    )
    return pd.to_numeric(s2, errors="coerce")

def clean_header_df(df: pd.DataFrame) -> pd.DataFrame:
    looks_like_placeholder = (
        any(isinstance(c, str) and c.startswith("KPI_") for c in df.columns) or
        any(isinstance(c, str) and c.startswith("Unnamed") for c in df.columns)
    )
    if looks_like_placeholder:
        new_cols = list(df.iloc[0].values)
        df = df.iloc[1:].copy()
        df.columns = new_cols

    df.columns = [str(c).strip() for c in df.columns]
    if "Month" not in df.columns:
        df.rename(columns={df.columns[0]: "Month"}, inplace=True)

    df["Month"] = pd.to_datetime(df["Month"], errors="coerce")

    for c in df.columns:
        if c == "Month":
            continue
        cname = c.lower()
        is_time = ("time" in cname) or ("hour" in cname) or ("min" in cname)
        df[c] = _numeric(df[c], treat_as_time=is_time)

    df = df.dropna(how="all").sort_values("Month").reset_index(drop=True)
    df["Year"] = df["Month"].dt.year
    df["MonthNum"] = df["Month"].dt.month
    df["MonthName"] = df["Month"].dt.strftime("%b")
    return df

# -------- Waste Breakup helpers --------
WB_TOTAL_COL_CANDIDATES = [
    "Total Waste (Tons)", "Total Waste Tons", "Total Waste", "Waste Total (Tons)",
]
STANDARD_KPI_COLS = {
    "Running Days", "Number of Accidents", "Production (Tons)", "Waste (Tons)",
    "Grade Changes", "Micron Changes", "Q2 (Tons)", "Reclaim (Tons)",
    "OEE %", "Waste %", "Break Recovery Time(Mint) / Break (min)",
    "Total Down Time (Hours)"
}
def detect_total_waste_col(cols: List[str]) -> Optional[str]:
    for c in cols:
        if c.strip().lower() == "total waste (tons)":
            return c
    for c in cols:
        cl = c.lower().strip()
        if any(k.lower() == cl for k in WB_TOTAL_COL_CANDIDATES):
            return c
        if ("total" in cl and "waste" in cl) and ("ton" in cl or "t" in cl):
            return c
    return None

def clean_waste_breakup_df(df: pd.DataFrame) -> pd.DataFrame:
    df = clean_header_df(df)
    num_cols = [c for c in df.columns if c not in {"Month","Year","MonthNum","MonthName"} and pd.api.types.is_numeric_dtype(df[c])]
    num_cols = [c for c in num_cols if c not in STANDARD_KPI_COLS]
    total_col = detect_total_waste_col(df.columns)
    if total_col and total_col not in num_cols:
        num_cols.append(total_col)
    keep = ["Month","Year","MonthNum","MonthName"] + num_cols
    df2 = df[keep].copy()
    df2.columns = [re.sub(r"\s+", " ", c).strip() for c in df2.columns]
    return df2

def is_waste_breakup_shape(df: pd.DataFrame) -> bool:
    if df is None or df.empty: return False
    has_month = "Month" in df.columns or (isinstance(df.columns[0], str))
    if not has_month: return False
    total_col = detect_total_waste_col(df.columns)
    numeric_candidates = [c for c in df.columns if c not in {"Month","Year","MonthNum","MonthName"} and pd.api.types.is_numeric_dtype(df[c])]
    if total_col and len(numeric_candidates) >= 2:
        return True
    non_kpi_numeric = [c for c in numeric_candidates if c not in STANDARD_KPI_COLS]
    return len(non_kpi_numeric) >= 3

# ====== Downtime detection helper ======
def is_downtime_breakup_shape(df: pd.DataFrame) -> bool:
    if df is None or df.empty: return False
    cols = {str(c).strip() for c in df.columns}
    if "Total Down Time (Hours)" in cols:
        return True
    dt_keys = {
        "No Job order","Planned Preventive Maintenance+Filter Change","Process",
        "Planning (G.C/M.C)","Mech. Plant","Mech. Utilities","Electrical Plant",
        "Electrical Utilities","Power House4","Rental Power"
    }
    return len(cols.intersection(dt_keys)) >= 3

# ------------------------------
# Year & Month pickers (per-page) — safer drop-in
# ------------------------------
import re
import hashlib

def _slug_key(prefix: str, label_prefix: str, salt: str = "") -> str:
    """
    Build a unique, stable Streamlit widget key from label_prefix.
    Avoids collisions across pages with similar labels.
    """
    slug = re.sub(r"[^a-z0-9]+", "_", (label_prefix or "").lower()).strip("_")
    h = hashlib.sha1((label_prefix + "|" + salt).encode("utf-8")).hexdigest()[:6]
    return f"{prefix}_{slug}_{h}" if slug else f"{prefix}_{h}"

def pick_year_month(
    df_all: pd.DataFrame,
    label_prefix: str = "",
    key_ns: str | None = None,   # <-- keep this param
) -> tuple[int | None, int | None, pd.DataFrame]:
    """
    Shows a Year and Month picker (Month optional) and returns:
      (selected_year, selected_monthnum, df_filtered_to_year_and_upto_month)
    """
    if df_all is None or df_all.empty or "Year" not in df_all.columns:
        return None, None, df_all

    # Ensure MonthNum is numeric (defensive)
    if "MonthNum" in df_all.columns:
        try:
            df_all = df_all.copy()
            df_all["MonthNum"] = pd.to_numeric(df_all["MonthNum"], errors="coerce")
        except Exception:
            pass

    # Build a unique key source so widget keys never collide on the same page
    key_source = f"{label_prefix}::{key_ns}" if key_ns else label_prefix

    years = sorted(df_all["Year"].dropna().astype(int).unique().tolist())
    if not years:
        return None, None, df_all
    default_year = years[-1]

    c1, c2 = st.columns([1, 1])

    with c1:
        sel_year = st.selectbox(
            f"{label_prefix}Year",
            options=years,
            index=(years.index(default_year) if default_year in years else 0),
            key=_slug_key("ym_year", key_source, salt="year"),
        )

    df_year = df_all[df_all["Year"] == sel_year].copy().sort_values("MonthNum")

    # Build month display map (supports missing names gracefully)
    months_df = (
        df_year[["MonthNum", "MonthName"]]
        .dropna(subset=["MonthNum"])
        .drop_duplicates()
        .sort_values("MonthNum")
    )
    months_df["MonthNum"] = months_df["MonthNum"].astype(int)
    month_name_by_num = {
        int(row["MonthNum"]): (
            str(row.get("MonthName")) if pd.notna(row.get("MonthName")) else f"M{int(row['MonthNum']):02d}"
        )
        for _, row in months_df.iterrows()
    }

    with c2:
        month_options = ["All months"] + [
            f"{m:02d} — {month_name_by_num.get(m, '')}" for m in months_df["MonthNum"].tolist()
        ]
        sel_month_label = st.selectbox(
            f"{label_prefix}Month",
            options=month_options,
            index=0,
            key=_slug_key("ym_month", key_source, salt="month"),
        )

    if sel_month_label == "All months" or not months_df.shape[0]:
        return int(sel_year), None, df_year

    # Parse the "MM — Name" label safely (support '—' or '-')
    try:
        raw = sel_month_label.split("—")[0] if "—" in sel_month_label else sel_month_label.split("-")[0]
        sel_monthnum = int(str(raw).strip())
    except Exception:
        # Fallback: use latest month in the year
        sel_monthnum = int(months_df["MonthNum"].max())

    df_trim = df_year[df_year["MonthNum"].astype(int) <= sel_monthnum].copy()
    return int(sel_year), int(sel_monthnum), df_trim

# =========================
# KPI helpers
# =========================
def kpi_delta(df_year: pd.DataFrame, col: str) -> Dict[str, Optional[float]]:
    if col not in df_year.columns:
        return {"value": None, "delta": None}
    d = df_year[["MonthNum", col]].dropna().sort_values("MonthNum")
    if d.empty:
        return {"value": None, "delta": None}
    cur = d.iloc[-1][col]
    prev = d.iloc[-2][col] if len(d) > 1 else np.nan
    return {"value": cur, "delta": (cur - prev) if pd.notna(prev) else None}

def kpi_delta_at_month(df_year: pd.DataFrame, col: str, sel_monthnum: int) -> Dict[str, Optional[float]]:
    if col not in df_year.columns:
        return {"value": None, "delta": None}
    d = df_year[["MonthNum", col]].dropna().sort_values("MonthNum")
    if d.empty or sel_monthnum not in d["MonthNum"].values:
        return {"value": None, "delta": None}
    cur = float(d.loc[d["MonthNum"] == sel_monthnum, col].iloc[0])
    prev_rows = d[d["MonthNum"] < sel_monthnum]
    prev = float(prev_rows.iloc[-1][col]) if len(prev_rows) > 0 else np.nan
    return {"value": cur, "delta": (cur - prev) if pd.notna(prev) else None}

def fmt_num(x, pct=False):
    if x is None or pd.isna(x):
        return "—"
    if pct:
        return f"{x:.2f}%"
    if isinstance(x, (int, np.integer)) or (isinstance(x, (float, np.floating)) and abs(x) >= 1000):
        return f"{x:,.0f}"
    return f"{x:.2f}"

def _apply_theme(fig: go.Figure) -> go.Figure:
    fig.update_layout(
        template="plotly_dark",
        paper_bgcolor="rgba(0,0,0,0)",
        plot_bgcolor="rgba(0,0,0,0)",
        font=dict(color=TEXT),
        margin=dict(l=40, r=24, t=60, b=40),
        colorway=COLORWAY,
    )
    return fig

def month_order_layout() -> Dict:
    return dict(
        categoryorder="array",
        categoryarray=["Jan", "Feb", "Mar", "Apr", "May", "Jun",
                       "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
    )

def line_fig(df_year: pd.DataFrame, y: str, title: str, color: str, pct=False) -> go.Figure:
    order = month_order_layout()
    dfp = df_year.sort_values("MonthNum")
    fig = px.line(dfp, x="MonthName", y=y, markers=True, title=title,
                  color_discrete_sequence=[color])
    if pct:
        fig.update_yaxes(ticksuffix=" %")
    fig.update_traces(mode="lines+markers", line=dict(width=3))
    fig.update_layout(xaxis_title="Month", yaxis_title=y, xaxis=order, hovermode="x unified")
    return _apply_theme(fig)

def area_fig(df_year: pd.DataFrame, y: str, title: str, color: str, pct=False) -> go.Figure:
    order = month_order_layout()
    dfp = df_year.sort_values("MonthNum")
    fig = px.area(dfp, x="MonthName", y=y, title=title,
                  color_discrete_sequence=[color])
    if pct:
        fig.update_yaxes(ticksuffix=" %")
    fig.update_traces(line=dict(width=2))
    fig.update_layout(xaxis_title="Month", yaxis_title=y, xaxis=order)
    return _apply_theme(fig)

def bar_fig(df_year: pd.DataFrame, y: str, title: str, color: str) -> go.Figure:
    order = month_order_layout()
    dfp = df_year.sort_values("MonthNum")
    fig = px.bar(dfp, x="MonthName", y=y, title=title,
                 color_discrete_sequence=[color])
    fig.update_traces(marker_line_width=0.5, marker_line_color="rgba(255,255,255,.3)")
    fig.update_layout(xaxis_title="Month", yaxis_title=y, xaxis=order, barmode="group")
    return _apply_theme(fig)

def stacked_bar_fig(df_year: pd.DataFrame, y1: str, y2: str, title: str, c1: str, c2: str) -> go.Figure:
    order = month_order_layout()
    dfp = df_year.sort_values("MonthNum").copy()
    fig = go.Figure()
    fig.add_bar(name=y1, x=dfp["MonthName"], y=dfp[y1], marker_color="#60A5FA")
    fig.add_bar(name=y2, x=dfp["MonthName"], y=dfp[y2], marker_color="#F59E0B")
    fig.update_layout(title=title, barmode="stack", xaxis_title="Month", yaxis_title="Tons", xaxis=order)
    return _apply_theme(fig)

def lollipop_fig(df_year: pd.DataFrame, y: str, title: str, color: str) -> go.Figure:
    order = month_order_layout()
    dfp = df_year.sort_values("MonthNum").copy()
    fig = go.Figure()
    fig.add_bar(x=dfp["MonthName"], y=dfp[y],
                marker_color=color, opacity=0.35, name=y)
    fig.add_scatter(x=dfp["MonthName"], y=dfp[y],
                    mode="markers", name=f"{y} points",
                    marker=dict(size=12, color=color, line=dict(color="white", width=1.4)))
    fig.update_layout(title=title, xaxis_title="Month", yaxis_title=y, xaxis=order, barmode="overlay")
    return _apply_theme(fig)

# === Unified export helpers (CSV / PDF / PPT) ===
def _fig_to_png(fig: go.Figure, width=1600, scale=2) -> Optional[bytes]:
    try:
        import plotly.io as pio  # requires kaleido
        png = fig.to_image(format="png", width=width, scale=scale)
        return png
    except Exception as e:
        st.info(f"PNG export needs 'kaleido'. Install it to enable image/PDF/PPT exports. Details: {e}")
        return None

def _build_pdf(figs: List[go.Figure], width=1600) -> Tuple[Optional[bytes], Optional[str]]:
    """
    PDF theme: dark background, SARKS logo on each page, footer bar with SARKS + clickable SARKS.
    One chart per page (full page). If multiple charts, multiple pages.
    """
    try:
        from PIL import Image
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.utils import ImageReader

        if not figs:
            return None, "No charts to export."

        # Render all figs to images
        imgs = []
        for f in figs:
            png = _fig_to_png(f, width=width, scale=2)
            if not png: 
                continue
            im = Image.open(BytesIO(png)).convert("RGB")
            imgs.append(im)
        if not imgs:
            return None, "PNG export failed (kaleido?)."

        logo_bytes = _logo_bytes()
        logo_reader = ImageReader(BytesIO(logo_bytes)) if logo_bytes else None

        pw, ph = A4  # portrait
        margin = 36  # 0.5 inch
        header_h = 54
        footer_h = 42

        buf = BytesIO()
        c = canvas.Canvas(buf, pagesize=(pw, ph))

        # Colors
        def rgb(hexstr):
            hexstr = hexstr.strip("#")
            return tuple(int(hexstr[i:i+2],16)/255 for i in (0,2,4))

        # Loop pages
        for im in imgs:
            # Background
            c.setFillColorRGB(*rgb(SLATE_900))
            c.rect(0, 0, pw, ph, fill=1, stroke=0)

            # Header area with logo
            if logo_reader:
                logo_w = 110
                logo_h = 36
                c.drawImage(logo_reader, margin, ph - margin - logo_h, width=logo_w, height=logo_h, mask='auto')

            # Title bar accent
            c.setFillColorRGB(*rgb(BRAND))
            c.rect(margin, ph - margin - header_h, pw - 2*margin, 2, fill=1, stroke=0)

            # Chart area sizing
            avail_w = pw - 2*margin
            avail_h = ph - (header_h + footer_h + 2*margin)
            # Fit image to area (keep aspect)
            img_w, img_h = im.size
            scale = min(avail_w / img_w, avail_h / img_h)
            draw_w, draw_h = img_w * scale, img_h * scale
            left = margin + (avail_w - draw_w)/2
            bottom = margin + footer_h + (avail_h - draw_h)/2
            c.drawImage(ImageReader(im), left, bottom, width=draw_w, height=draw_h)

            # Footer bar
            c.setFillColorRGB(*rgb(SLATE_760))
            c.rect(0, 0, pw, footer_h, fill=1, stroke=0)
            c.setFillColorRGB(1,1,1)
            c.setFont("Helvetica", 9.5)
            c.drawString(margin, 14, FOOTER_LEFT)
            # Right side: "developed by SARKS" with link on SARKS
            right_text = FOOTER_RIGHT_TEXT + FOOTER_RIGHT_BRAND
            tw = c.stringWidth(right_text, "Helvetica", 9.5)
            x_right = pw - margin - tw
            c.drawString(x_right, 14, right_text)
            # Link only on SARKS word
            sarks_w = c.stringWidth(FOOTER_RIGHT_BRAND, "Helvetica", 9.5)
            sarks_x = x_right + c.stringWidth(FOOTER_RIGHT_TEXT, "Helvetica", 9.5)
            c.linkURL(DEV_LINK, (sarks_x, 8, sarks_x + sarks_w, 20), relative=0)
            c.showPage()

        c.save()
        return buf.getvalue(), None
    except Exception as e:
        return None, f"PDF export requires kaleido + Pillow + ReportLab. Details: {e}"

def _build_ppt(figs: List[go.Figure], title: str, narrative_lines: List[str]) -> Tuple[Optional[bytes], Optional[str]]:
    """
    PPT theme: dark background + SARKS logo on every slide (no footer bar).
    Slides:
      • Title slide
      • Executive summary (optional)
      • One slide per chart (with explanation box)
      • Conclusion slide
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from io import BytesIO
        import os

        if not figs:
            return None, "No charts to export."

        prs = Presentation()

        # --- Load logo safely ---
        def _logo_bytes() -> bytes:
            candidates = [
                "SARKS-01.png",
                "./assets/SARKS-01.png",
                "./static/SARKS-01.png"
            ]
            for path in candidates:
                if os.path.exists(path):
                    with open(path, "rb") as f:
                        return f.read()
            return b""  # always return bytes, never None

        logo_bytes = _logo_bytes()

        # --- theme colors ---
        DARK_BG = RGBColor(0x12, 0x16, 0x1D)  # SLATE_900
        WHITE   = RGBColor(0xEA, 0xEF, 0xF5)
        BRANDC  = RGBColor(0xF5, 0x60, 0x23)  # SARKS orange

        def set_dark_bg(slide):
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = DARK_BG

        def add_logo(slide, height_in=0.6):
            """Safely add logo if available."""
            if logo_bytes and len(logo_bytes) > 0:
                slide.shapes.add_picture(BytesIO(logo_bytes), Inches(0.25), Inches(0.25), height=Inches(height_in))

        # --- Title slide ---
        title_slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        set_dark_bg(title_slide)
        add_logo(title_slide)

        tb = title_slide.shapes.add_textbox(Inches(0.25), Inches(1.2), Inches(9.5), Inches(1.0))
        tf = tb.text_frame; tf.clear()
        p = tf.paragraphs[0]
        # SAFE: title may be None
        p.text = (title or "KPI Dashboard")
        p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = WHITE

        sub_tb = title_slide.shapes.add_textbox(Inches(0.25), Inches(2.1), Inches(9.5), Inches(0.8))
        stf = sub_tb.text_frame; stf.clear()
        ps = stf.paragraphs[0]
        ps.text = "BOPET FilmLine KPI Dashboard • SARKS IntelliLine"
        ps.font.size = Pt(18); ps.font.color.rgb = BRANDC

        # --- Executive summary ---
        # SAFE: narrative_lines may contain None values
        safe_lines = [str(l) for l in (narrative_lines or []) if l]
        if safe_lines:
            s = prs.slides.add_slide(prs.slide_layouts[6])
            set_dark_bg(s); add_logo(s)
            tb = s.shapes.add_textbox(Inches(0.25), Inches(0.85), Inches(9.5), Inches(0.6))
            tf = tb.text_frame; tf.clear()
            p = tf.paragraphs[0]
            p.text = "Executive Summary"
            p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = WHITE

            box = s.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9.0), Inches(5.5))
            tfb = box.text_frame; tfb.clear()
            for i, line in enumerate(safe_lines):
                para = tfb.add_paragraph() if i > 0 else tfb.paragraphs[0]
                para.text = line
                para.font.size = Pt(14); para.font.color.rgb = WHITE

        # --- One chart per slide ---
        for i, fig in enumerate(figs, start=1):
            png = _fig_to_png(fig, width=1400, scale=2)
            if not png:
                continue
            s = prs.slides.add_slide(prs.slide_layouts[6])
            set_dark_bg(s); add_logo(s)

            # Chart title
            # SAFE: fig.layout.title may exist with .text == None
            t_txt = (getattr(getattr(fig.layout, "title", None), "text", None) or f"Chart {i}")
            tb = s.shapes.add_textbox(Inches(0.25), Inches(0.85), Inches(9.5), Inches(0.6))
            tf = tb.text_frame; tf.clear()
            p = tf.paragraphs[0]
            p.text = t_txt
            p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = WHITE

            # Chart image
            s.shapes.add_picture(BytesIO(png), Inches(0.5), Inches(1.6), width=Inches(9.0))

            # Explanation box
            exp = s.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9.0), Inches(1.0))
            etf = exp.text_frame; etf.clear()
            head = etf.paragraphs[0]
            head.text = "Detailed explanation"
            head.font.bold = True; head.font.size = Pt(12); head.font.color.rgb = WHITE

            default_bullets = [
                "Monthly trend with SARKS dark theme.",
                "Spot peaks/dips and seasonal effects.",
                "Compare vs. last month and YTD.",
                "Drive corrective actions for production, waste, OEE, downtime."
            ]
            for b in default_bullets:
                para = etf.add_paragraph()
                para.text = f"• {b}"
                para.font.size = Pt(11); para.font.color.rgb = WHITE

        # --- Conclusion ---
        concl = prs.slides.add_slide(prs.slide_layouts[6])
        set_dark_bg(concl); add_logo(concl)
        tb = concl.shapes.add_textbox(Inches(0.25), Inches(0.85), Inches(9.5), Inches(0.6))
        tf = tb.text_frame; tf.clear()
        p = tf.paragraphs[0]
        p.text = "Conclusion & Next Steps"
        p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = WHITE

        box = concl.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9.0), Inches(5.5))
        tfb = box.text_frame; tfb.clear()
        for txt in [
            "Focus on reducing top contributors in Waste & Downtime paretos.",
            "Sustain OEE by stabilizing changeovers and recovery.",
            "Track forecasts vs. actuals next month and adjust plan."
        ]:
            para = tfb.add_paragraph() if tfb.paragraphs[0].text else tfb.paragraphs[0]
            para.text = txt
            para.font.size = Pt(14); para.font.color.rgb = WHITE

        out = BytesIO()
        prs.save(out)
        return out.getvalue(), None

    except Exception as e:
        return None, f"PPT export requires python-pptx (and kaleido for chart images). Details: {e}"


import os

def _logo_bytes() -> bytes:
    """
    Always return valid bytes.
    - Reads SARKS-01.png if available
    - Otherwise returns empty bytes (safe fallback)
    """
    candidates = [
        "SARKS-01.png",
        "./assets/SARKS-01.png",
        "./static/SARKS-01.png"
    ]
    for path in candidates:
        if os.path.exists(path):
            with open(path, "rb") as f:
                return f.read()
    return b""   # never return None

def _csv_bytes(df: pd.DataFrame) -> bytes:
    return df.to_csv(index=False).encode("utf-8")

def _download_row(tab_name: str, csv_df: Optional[pd.DataFrame], figs: List[go.Figure], narrative_lines: List[str]):
    st.markdown("---")
    st.markdown("### Downloads")
    col1, col2, col3 = st.columns(3)
    with col1:
        if csv_df is not None and not csv_df.empty:
            st.download_button(
                f"⬇ CSV — {tab_name}",
                data=_csv_bytes(csv_df),
                file_name=f"{tab_name.lower().replace(' ','_')}_{datetime.now():%Y%m%d}.csv",
                use_container_width=True
            )
        else:
            st.caption("No table for CSV.")
    with col2:
        pdf_bytes, pdf_err = _build_pdf(figs)
        if pdf_bytes:
            st.download_button(
                f"⬇ PDF — {tab_name}",
                data=pdf_bytes,
                file_name=f"{tab_name.lower().replace(' ','_')}_{datetime.now():%Y%m%d}.pdf",
                mime="application/pdf",
                use_container_width=True
            )
        else:
            st.info(pdf_err or "No charts for PDF.")
    with col3:
        ppt_bytes, ppt_err = _build_ppt(figs, f"{tab_name} — SARKS IntelliLine", narrative_lines)
        if ppt_bytes:
            st.download_button(
                f"⬇ PPT — {tab_name}",
                data=ppt_bytes,
                file_name=f"{tab_name.lower().replace(' ','_')}_{datetime.now():%Y%m%d}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True
            )
        else:
            st.info(ppt_err or "No charts for PPT.")

# NEW: add forecast range marker utility
def add_range_marker(fig: go.Figure, x_label: str, lo: float, mid: float, hi: float, name="Forecast") -> go.Figure:
    fig.add_scatter(
        x=[x_label], y=[mid],
        mode="markers",
        name=name,
        marker=dict(symbol="diamond", size=12, line=dict(color="white", width=1.4)),
        error_y=dict(type="data", array=[max(0.0, hi - mid)], arrayminus=[max(0.0, mid - lo)])
    )
    return fig

# NEW: generic multi-compare function (multi-color, dual-axis when needed)
def _is_pct_col(name: str) -> bool:
    n = name.strip().lower()
    return n.endswith("%") or "percent" in n or n.endswith(" pct")

def compare_multi_kpis_fig(df_year: pd.DataFrame, base: str, others: List[str]) -> Optional[go.Figure]:
    if base not in df_year.columns:
        return None
    dfp = df_year.sort_values("MonthNum").copy()
    order = month_order_layout()
    fig = go.Figure()
    base_is_pct = _is_pct_col(base)
    base_axis = "y2" if base_is_pct else "y"
    base_color = "#60A5FA"
    fig.add_bar(name=base, x=dfp["MonthName"], y=dfp[base],
                marker_color=base_color, opacity=0.85, yaxis=base_axis)

    color_pool = COLORWAY[2:] + COLORWAY[:2]
    for idx, o in enumerate(others):
        if o not in dfp.columns:
            continue
        o_is_pct = _is_pct_col(o)
        axis = "y2" if o_is_pct and not base_is_pct else ("y" if not o_is_pct and base_is_pct else base_axis)
        fig.add_scatter(name=o, x=dfp["MonthName"], y=dfp[o],
                        mode="lines+markers",
                        marker=dict(size=8),
                        line=dict(width=3, color=color_pool[idx % len(color_pool)]),
                        yaxis=axis)

    if base_is_pct:
        fig.update_layout(
            yaxis=dict(title="Value"),
            yaxis2=dict(title=base, overlaying="y", side="right", ticksuffix=" %")
        )
    else:
        fig.update_layout(
            yaxis=dict(title=base),
            yaxis2=dict(title="Percent KPIs", overlaying="y", side="right", ticksuffix=" %")
        )

    fig.update_layout(
        title=f"{base} — Comparison",
        xaxis_title="Month", xaxis=month_order_layout(), hovermode="x unified",
        legend=dict(orientation="h", yanchor="bottom", y=1.02, x=0)
    )
    return _apply_theme(fig)

def corr_heatmap(df: pd.DataFrame, title: str) -> Optional[go.Figure]:
    num = df.select_dtypes(include=[np.number]).drop(columns=["Year", "MonthNum"], errors="ignore")
    if num.shape[1] < 2:
        return None
    corr = num.corr().abs()
    fig = go.Figure(data=go.Heatmap(
        z=corr.values,
        x=list(corr.columns),
        y=list(corr.index),
        zmin=0, zmax=1,
        colorscale="Blues"
    ))
    fig.update_layout(title=title)
    return _apply_theme(fig)

# Aggregations
SUM_KPIS = {"Running Days", "Number of Accidents", "Production (Tons)", "Waste (Tons)",
            "Grade Changes", "Micron Changes", "Q2 (Tons)", "Reclaim (Tons)"}
AVG_KPIS = {"OEE %", "Waste %", "Break Recovery Time(Mint) / Break (min)", "Total Down Time (Hours)"}

def yearly_value(df_year: pd.DataFrame, col: str) -> Optional[float]:
    if col not in df_year.columns:
        return None
    s = pd.to_numeric(df_year[col], errors="coerce").dropna()
    if s.empty:
        return None
    if col in SUM_KPIS:
        return float(s.sum())
    if col in AVG_KPIS:
        return float(s.mean())
    return float(s.sum())

def waste_composition_fixed(df_year: pd.DataFrame) -> Tuple[Optional[pd.DataFrame], Optional[str], Optional[float]]:
    need = {"Production (Tons)", "Waste (Tons)"}
    if not need.issubset(df_year.columns):
        return None, "Need 'Production (Tons)' and 'Waste (Tons)' for the donut.", None
    prod_sum = float(pd.to_numeric(df_year["Production (Tons)"], errors="coerce").fillna(0).sum())
    waste_sum = float(pd.to_numeric(df_year["Waste (Tons)"], errors="coerce").fillna(0).sum())
    total = prod_sum + waste_sum
    if total <= 0:
        return None, "Total production is zero; cannot compute composition.", None
    donut = pd.DataFrame({"Category": ["Good Production", "Waste"], "Tons": [prod_sum, waste_sum]})
    waste_pct_of_total = (waste_sum / total) * 100.0
    return donut, None, waste_pct_of_total

# =========================
# Styles (header & cards) + Footer + Modal
# =========================
LOGO64 = _logo_b64()
st.markdown(
    f"""
    <style>
      .stApp {{
        background: radial-gradient(1200px 700px at 12% 0%, {SLATE_900} 0%, {SLATE_970} 70%) fixed;
      }}
      .stApp::before {{
        content: "";
        position: fixed; inset: 0;
        background: url('data:image/png;base64,{LOGO64 if LOGO64 else ""}') no-repeat center 18%;
        background-size: 62% auto;
        opacity: .06;
        pointer-events: none;
      }}
      .block-container {{ max-width: 1540px; padding-top: .2rem; padding-bottom: 90px; }}
      section[data-testid="stSidebar"] {{
        background: linear-gradient(180deg, {SLATE_900}, {SLATE_830});
        border-right: 1px solid rgba(255,255,255,.06);
        position: relative; overflow: hidden;
      }}
      section[data-testid="stSidebar"]::before {{
        content: "";
        position: absolute; inset: 0 0 -20% 0;
        background: url('data:image/png;base64,{LOGO64 if LOGO64 else ""}') no-repeat center 80%;
        background-size: 120% auto;
        opacity: .06; pointer-events: none; filter: saturate(0.9);
      }}
      section[data-testid="stSidebar"] .stSelectbox label,
      section[data-testid="stSidebar"] .stCheckbox label,
      section[data-testid="stSidebar"] .stRadio label,
      section[data-testid="stSidebar"] .stSlider label {{
        color: {BRAND} !important; font-weight: 700;
      }}
      section[data-testid="stSidebar"] .stButton>button {{
        background: {BRAND}; color: white; border: 0;
        box-shadow: 0 4px 14px rgba(245,96,35,.35);
      }}
      section[data-testid="stSidebar"] .stCheckbox>div>div {{ border-color: {BRAND} !important; }}
      section[data-testid="stSidebar"] .stSelectbox [data-baseweb="select"]>div {{
        border-color: rgba(245,96,35,.55);
      }}
      .ppk-header {{
          display: grid; grid-template-columns: {HEADER_LOGO_COL_PX}px 1fr;
          align-items: center; gap: 18px;
          background: linear-gradient(180deg, {SLATE_900}, {SLATE_830});
          border: 1px solid rgba(255,255,255,0.06);
          border-radius: 16px; padding: 8px 14px;
          box-shadow: 0 8px 24px rgba(0,0,0,.35);
          min-height: {HEADER_MIN_H}px;
          position: sticky; top: 0; z-index: 50;
      }}
      .ppk-logo {{ height: {LOGO_HEIGHT_PX}px; max-height:{LOGO_HEIGHT_PX}px; object-fit: contain; width:100%; }}
      .ppk-title {{ color:{TEXT}; font-size:1.55rem; font-weight:900; letter-spacing:.3px; line-height:1.2; }}
      .ppk-sub   {{ color:{MUTED}; font-size:.92rem; margin-top:2px; }}
      .metric-card {{
          position: relative; background:{SLATE_800};
          border:1px solid rgba(255,255,255,.07); border-radius:18px;
          padding:14px 16px 12px 16px; box-shadow:0 6px 18px rgba(0,0,0,.25);
      }}
      .metric-card:before {{
          content:""; position:absolute; left:0; top:0; bottom:0; width:6px;
          background:{BRAND}; border-top-left-radius:18px; border-bottom-left-radius:18px;
      }}
      .metric-caption {{ color:{MUTED}; font-size:.9rem; margin-bottom:4px; }}
      .big-value {{ color:{TEXT}; font-size:1.9rem; font-weight:800; line-height:1.15; }}
      .delta-badge {{ display:inline-block; padding:3px 10px; border-radius:999px; font-size:.85rem; font-weight:700;
                      background:{SLATE_720}; border:1px solid rgba(255,255,255,.12); color:{TEXT}; }}
      .year-card {{ background: linear-gradient(180deg, rgba(245,96,35,.12), rgba(245,96,35,.05));
                   border:1px solid rgba(245,96,35,.35); border-radius:16px; padding:14px 16px; }}
      .year-val {{ font-size:1.8rem; font-weight:900; color:{BRAND}; text-shadow:0 0 10px rgba(245,96,35,.25); }}
      .year-chip {{ background: rgba(245,96,35,.15); color:{BRAND}; border:1px solid rgba(245,96,35,.35);
                   padding:4px 10px; border-radius:999px; font-weight:700; font-size:.85rem; }}
      .ppk-panel {{ background:{SLATE_850}; border:1px solid rgba(255,255,255,.06);
                   border-radius:14px; padding:10px 14px 6px 14px; box-shadow:0 6px 16px rgba(0,0,0,.24); }}
      .ppk-titlebar {{ color:{TEXT}; font-weight:800; margin-top:4px; }}
      .ppk-accent {{ height:3px; width:120px; margin:4px 0 6px 0; background:{BRAND}; }}
      /* App footer bar (fixed) */
      .app-footer-bar {{
        position: fixed; left: 0; right: 0; bottom: 0;
        background: {SLATE_760};
        border-top: 1px solid rgba(255,255,255,.15);
        color: {TEXT};
        font-size: 14px;
        padding: 8px 14px;
        z-index: 100;
      }}
      .app-footer-bar .row {{
        display:flex; align-items:center; justify-content:space-between; gap: 8px;
        max-width: 1540px; margin: 0 auto;
      }}
      .sarks-link button {{
        background: none !important;
        border: none !important;
        padding: 0 !important;
        color: {TEXT} !important;
        font-weight: 700 !important;
        text-decoration: underline !important;
        cursor: pointer !important;
      }}
      /* Modal */
      .modal-overlay {{
        position: fixed; inset: 0;
        background: rgba(0,0,0,.55);
        display: flex; align-items: center; justify-content: center;
        z-index: 999;
      }}
      .modal-card {{
        width: min(780px, 92vw);
        background: {SLATE_830};
        border: 1px solid rgba(255,255,255,.12);
        border-radius: 16px;
        box-shadow: 0 24px 80px rgba(0,0,0,.5);
        padding: 18px 20px;
        color: {TEXT};
      }}
      .modal-card h3 {{ margin: 2px 0 8px 0; }}
      .modal-card .muted {{ color: {MUTED}; }}
      .modal-card .brand {{ color: {BRAND}; font-weight: 800; }}
      .modal-close button {{
        background: {BRAND} !important; color: white !important; border: 0 !important;
        padding: 6px 12px !important; border-radius: 10px !important;
      }}
    </style>
    """,
    unsafe_allow_html=True,
)

def render_footer_bar():
    st.markdown(
        f"""
        <style>
          .app-footer-bar {{
            position: fixed; left: 0; right: 0; bottom: 0;
            background:
              linear-gradient(90deg, rgba(245,96,35,.18) 0%, rgba(245,96,35,.10) 42%, rgba(245,96,35,0) 78%),
              linear-gradient(180deg, {SLATE_900}, {SLATE_830});
            border-top: 1px solid rgba(255,255,255,.15);
            color: {TEXT};
            font-size: 14px; padding: 12px 18px;
            z-index: 9999; box-shadow: 0 -6px 20px rgba(0,0,0,.28);
            font-weight: 700;
          }}
          .app-footer-bar .row {{
            display:flex; align-items:center; justify-content:space-between; gap: 10px;
            max-width: 1540px; margin: 0 auto;
          }}
          .block-container {{ padding-bottom: 96px; }}
        </style>

        <div class="app-footer-bar">
          <div class="row">
            <div style="text-align:left;">
              © 2025 • SARKS IntelliLine — BOPET FilmLine KPI Dashboard
            </div>
            <div style="text-align:right;">
              Version-01
            </div>
          </div>
        </div>
        """,
        unsafe_allow_html=True,
    )

def render_about_developer_page():
    """
    Professional About page.
    Renders with components.html so HTML never shows as plain text.
    Neutral text; subtle brand gradient in the hero.
    """
    import textwrap
    from streamlit.components.v1 import html as st_html

    css = textwrap.dedent(f"""
    <style>
      :root {{
        --brand: {BRAND};
        --slate900: {SLATE_900};
        --slate830: {SLATE_830};
        --text: {TEXT};
        --muted: {MUTED};
      }}

      .about-wrap {{
        max-width: 1100px;
        margin: 10px auto 24px auto;
        padding: 0 8px;
        color: var(--text);
        font-family: system-ui, -apple-system, Segoe UI, Roboto, Arial, sans-serif;
      }}

      /* HERO */
      .about-hero {{
        position: relative;
        background:
          radial-gradient(1100px 650px at -10% -20%, rgba(245,96,35,.18) 0%, rgba(245,96,35,.08) 42%, rgba(245,96,35,0) 78%),
          linear-gradient(180deg, var(--slate900), var(--slate830));
        border: 1px solid rgba(255,255,255,.10);
        border-radius: 18px;
        padding: 18px 20px 20px 20px;
        box-shadow: 0 10px 28px rgba(0,0,0,.35);
        overflow: hidden;
      }}
      .about-hero::after {{
        content:"";
        position:absolute; inset:0;
        background: linear-gradient(90deg, rgba(245,96,35,.20), rgba(245,96,35,0));
        opacity:.22; pointer-events:none;
      }}

      .about-title {{
        font-weight: 900; line-height: 1.15;
        font-size: 2rem; letter-spacing: .2px; margin: 2px 0 2px 0;
      }}
      .about-sub {{
        color: var(--muted); font-size: .98rem; margin-bottom: 8px;
      }}

      .badge-row {{ display:flex; gap:10px; flex-wrap:wrap; margin: 6px 0 2px 0; }}
      .badge, .badge-muted {{
        display:inline-block; padding:7px 12px; border-radius:999px;
        background: rgba(255,255,255,.06);
        border:1px solid rgba(255,255,255,.14);
        color: var(--text); font-weight:800; font-size:.85rem;
      }}
      .badge {{ border-color: rgba(245,96,35,.45); background: rgba(245,96,35,.12); }}

      .link-row {{ margin-top: 10px; }}
      .link-out {{
        color: var(--text); text-decoration:none; font-weight:800;
        border-bottom:1px dashed var(--brand);
      }}
      .link-out:hover {{ color:#fff; border-bottom-color:#fff; }}

      /* GRID */
      .about-grid {{
        margin-top: 18px;
        display: grid;
        grid-template-columns: 1fr 1fr;
        gap: 14px;
      }}
      @media (max-width: 900px) {{ .about-grid {{ grid-template-columns: 1fr; }} }}

      .card {{
        background: var(--slate830);
        border: 1px solid rgba(255,255,255,.10);
        border-radius: 16px;
        padding: 18px;
        box-shadow: 0 12px 28px rgba(0,0,0,.35);
      }}
      .card h3 {{ margin: 0 0 8px 0; font-size: 1.06rem; font-weight: 800; }}
      .card p {{ margin:.35rem 0; opacity:.95; }}
      .card ul {{ margin:.2rem 0 0 1.2rem; line-height:1.6; }}
      .card li {{ margin:.15rem 0; }}

      /* Actions */
      .about-actions {{ margin-top: 16px; display:flex; gap:10px; flex-wrap:wrap; }}
      .btn-back {{
        display:inline-block; padding:10px 16px; border-radius:10px;
        background: var(--brand); color:white; font-weight:800; text-decoration:none;
        border:0; box-shadow:0 6px 16px rgba(245,96,35,.35);
      }}
      .btn-back:hover {{ filter: brightness(1.05); }}
    </style>
    """)

    html = textwrap.dedent(f"""
    {css}
    <div class="about-wrap">
      <section class="about-hero">
        <div class="about-title">About the Developer</div>
        <div class="about-sub">SARKS • IntelliLine • BOPET FilmLine KPI Dashboard</div>

        <div class="badge-row">
          <span class="badge">Sheikh Abdul Rehman bin Khalid Sharif — SARKS</span>
          <span class="badge-muted">Developer • Data Engineer</span>
        </div>

        <div class="link-row">
          <a class="link-out" href="{DEV_LINK}" target="_blank" rel="noopener">LinkedIn • Sheikh Abdul Rehman</a>
        </div>
      </section>

      <section class="about-grid">
        <div class="card">
          <h3>What this model does</h3>
          <ul>
            <li><b>KPI Dashboard</b> — visualizes production, waste, downtime, OEE, and quality metrics.</li>
            <li><b>Breakups</b> — in-depth analysis of waste and downtime using Pareto charts and factor shares.</li>
            <li><b>Forecasting</b> — next-month predictions with uncertainty estimation.</li>
            <li><b>Exports</b> — ready-to-use CSV, PDF, and PPT reports.</li>
          </ul>
        </div>

        <div class="card">
          <h3>Highlights</h3>
          <ul>
            <li>Enterprise-grade dark theme aligned with SARKS branding.</li>
            <li>Interactive analytics powered by Plotly.</li>
            <li>Automated validation checks to minimize reporting errors.</li>
          </ul>
        </div>

        <div class="card">
          <h3>Tech stack</h3>
          <p>Python • Streamlit • Pandas • NumPy • Plotly • ReportLab • python-pptx</p>
          <p>Design system: SARKS orange <code>{BRAND}</code> with slate backgrounds.</p>
        </div>

        <div class="card">
          <h3>Vision & Future Work</h3>
          <p>This model is built to provide transparency, efficiency, and data-driven decision making for film line KPIs. Future enhancements include integration with real-time IoT data sources and advanced machine learning–based forecasting.</p>
        </div>
      </section>

      <div class="about-actions">
        <a class="btn-back" href=".">← Back to Dashboard</a>
      </div>
    </div>
    """)

    # Render as HTML
    st_html(html, height=820, scrolling=True)


# =========================
# Session state
# =========================
# =========================
# Session state (init) + AboutDev one-shot
# =========================
active = st.session_state.get("active_page", "Home")

if "sarks_df" not in st.session_state:
    st.session_state.sarks_df = pd.DataFrame()
if "waste_breakup_df" not in st.session_state:
    st.session_state.waste_breakup_df = pd.DataFrame()
if "downtime_breakup_df" not in st.session_state:
    st.session_state.downtime_breakup_df = pd.DataFrame()
if "waste_raw_df" not in st.session_state:
    st.session_state.waste_raw_df = pd.DataFrame()
if "downtime_raw_df" not in st.session_state:
    st.session_state.downtime_raw_df = pd.DataFrame()

if "active_page" not in st.session_state:
    st.session_state.active_page = "Home"
if "breakup_mode" not in st.session_state:
    st.session_state.breakup_mode = "Waste"
if "show_dev_modal" not in st.session_state:
    st.session_state.show_dev_modal = False

# --- One-shot route: consume ?dev=1, then clear it and rerun
dev_qs = st.query_params.get("dev")
if dev_qs == "1":
    st.session_state.active_page = "AboutDev"
    st.query_params.clear()   # remove ?dev=1 so it doesn't stick
    st.rerun()

# If session says About page, render it and stop
if st.session_state.get("active_page") == "AboutDev":
    render_about_developer_page()
    st.stop()

# =========================
# Breakup expected columns
# =========================
WB_EXPECTED_COLUMNS = [
    "Month",
    "Monthly Waste (Tons)",
    "Process",
    "Planning(G.C/M.C)",
    "Special Films",
    "Electrical",
    "Mechanical",
    "Power House",
    "No Job Order",
    "Waste %",
]
DT_EXPECTED_COLUMNS = [
    "Month",
    "Total Down Time (Hours)",
    "No Job order",
    "Planned Preventive Maintenance+Filter Change",
    "Process",
    "Planning (G.C/M.C)",
    "Mech. Plant",
    "Mech. Utilities",
    "Electrical Plant",
    "Electrical Utilities",
    "Power House4",
    "Rental Power",
]

# =========================
# Sidebar THEME (SARKS gradient + brand colors)
# =========================
st.markdown(
    f"""
    <style>
      /* Sidebar container */
      section[data-testid="stSidebar"] {{
        background:
          linear-gradient(90deg, rgba(245,96,35,.18) 0%, rgba(245,96,35,.10) 42%, rgba(245,96,35,0) 78%),
          linear-gradient(180deg, {SLATE_900}, {SLATE_830});
        border-right: 1px solid rgba(255,255,255,.06);
        position: relative; 
        overflow: hidden;
      }}

      /* Subtle SARKS logo watermark inside sidebar */
      section[data-testid="stSidebar"]::before {{
        content: "";
        position: absolute; 
        inset: 0;
        background: url('data:image/png;base64,{LOGO64 if LOGO64 else ""}') no-repeat center 90%;
        background-size: 80% auto;
        opacity: .07; 
        pointer-events: none;
      }}

      /* Sidebar labels */
      section[data-testid="stSidebar"] .stSelectbox label,
      section[data-testid="stSidebar"] .stCheckbox label,
      section[data-testid="stSidebar"] .stRadio label,
      section[data-testid="stSidebar"] .stSlider label {{
        color: {BRAND} !important; 
        font-weight: 700;
      }}

      /* Buttons */
      section[data-testid="stSidebar"] .stButton>button {{
        background: {BRAND}; 
        color: white; 
        border: 0;
        border-radius: 10px;
        box-shadow: 0 4px 14px rgba(245,96,35,.35);
      }}

      /* Checkbox border brand */
      section[data-testid="stSidebar"] .stCheckbox>div>div {{ 
        border-color: {BRAND} !important; 
      }}

      /* Selectbox highlight */
      section[data-testid="stSidebar"] .stSelectbox [data-baseweb="select"]>div {{
        border-color: rgba(245,96,35,.55);
      }}
    </style>
    """,
    unsafe_allow_html=True
)

# =========================
# Sidebar CONTENT
# =========================
with st.sidebar:
    st.markdown("### Data")
    workbook_file = st.file_uploader("Upload workbook (Excel preferred)", 
                                     type=["xlsx", "xls", "csv"], 
                                     key="workbook_uploader")
    st.caption("Put main KPIs on a sheet (e.g., 'KPI'/'Summary'), Waste Breakup on another (e.g., 'Waste Breakup'), and Downtime Breakup on a third (e.g., 'Downtime').")

def _try_read_csv_or_excel_first_sheet(f):
    try:
        if f.name.lower().endswith(".csv"):
            return {"KPI": pd.read_csv(f)}, "KPI", None, None
        else:
            xl = pd.ExcelFile(f)
            return {sn: pd.read_excel(f, sheet_name=sn) for sn in xl.sheet_names}, None, None, None
    except Exception as e:
        st.error(f"Failed to read file: {e}")
        return {}, None, None, None

def _classify_sheets(sheet_map: Dict[str, pd.DataFrame]) -> Tuple[Optional[str], Optional[str], Optional[str], Dict[str, pd.DataFrame]]:
    kpi_sheet = waste_sheet = dt_sheet = None
    cleaned_preview = {}

    def _name_hint(name: str) -> str:
        n = name.lower()
        if "waste" in n: return "waste"
        if "down" in n or "dt" in n: return "dt"
        if any(k in n for k in ["kpi","summary","main","annual","dashboard","report"]): return "kpi"
        return ""

    for sn in sheet_map:
        hint = _name_hint(sn)
        if hint == "kpi" and kpi_sheet is None:
            kpi_sheet = sn
        elif hint == "waste" and waste_sheet is None:
            waste_sheet = sn
        elif hint == "dt" and dt_sheet is None:
            dt_sheet = sn

    for sn, raw in sheet_map.items():
        try:
            cleaned = clean_header_df(raw.copy())
            cleaned_preview[sn] = cleaned
        except Exception:
            continue

    if kpi_sheet is None:
        for sn, cdf in cleaned_preview.items():
            cols = set(cdf.columns)
            if {"Production (Tons)"}.issubset(cols):
                kpi_sheet = sn
                break

    if waste_sheet is None:
        for sn, cdf in cleaned_preview.items():
            if is_waste_breakup_shape(cdf):
                waste_sheet = sn
                break

    if dt_sheet is None:
        for sn, cdf in cleaned_preview.items():
            if is_downtime_breakup_shape(cdf):
                dt_sheet = sn
                break

    return kpi_sheet, waste_sheet, dt_sheet, cleaned_preview

# Load the workbook once selected
if workbook_file is not None:
    sheets, _, _, _ = _try_read_csv_or_excel_first_sheet(workbook_file)
    if sheets:
        kpi_sn, waste_sn, dt_sn, cleaned_preview = _classify_sheets(sheets)

        if kpi_sn and kpi_sn in sheets:
            st.session_state.sarks_df = clean_header_df(sheets[kpi_sn].copy())
        elif len(sheets) == 1:
            only_df = list(sheets.values())[0]
            st.session_state.sarks_df = clean_header_df(only_df.copy())

        if waste_sn and waste_sn in sheets:
            st.session_state.waste_raw_df = sheets[waste_sn].copy()
            try:
                st.session_state.waste_breakup_df = clean_waste_breakup_df(st.session_state.waste_raw_df.copy())
            except Exception:
                st.session_state.waste_breakup_df = clean_header_df(st.session_state.waste_raw_df.copy())
        else:
            st.session_state.waste_raw_df = pd.DataFrame()
            st.session_state.waste_breakup_df = pd.DataFrame()

        if dt_sn and dt_sn in sheets:
            st.session_state.downtime_raw_df = sheets[dt_sn].copy()
            try:
                st.session_state.downtime_breakup_df = clean_header_df(st.session_state.downtime_raw_df.copy())
            except Exception:
                st.session_state.downtime_breakup_df = pd.DataFrame()
        else:
            st.session_state.downtime_raw_df = pd.DataFrame()
            st.session_state.downtime_breakup_df = pd.DataFrame()

        with st.sidebar.expander("Detected sheets", expanded=False):
            st.write(f"• KPI sheet: **{kpi_sn or '—'}**")
            st.write(f"• Waste Breakup sheet: **{waste_sn or '—'}**")
            st.write(f"• Downtime Breakup sheet: **{dt_sn or '—'}**")

# Working copies
df = st.session_state.sarks_df.copy()
df_wb = st.session_state.waste_breakup_df.copy()
df_dt = st.session_state.downtime_breakup_df.copy()

# ---- Years available (safe construction)
years_all = sorted(set(
    (df["Year"].dropna().astype(int).tolist() if not df.empty else [])
    + (df_wb["Year"].dropna().astype(int).tolist() if not df_wb.empty else [])
    + (df_dt["Year"].dropna().astype(int).tolist() if not df_dt.empty else [])
))
year_default = max(years_all) if years_all else datetime.now().year
year = year_default  # keep variable for exports

df_year_default = df[df["Year"] == year_default].copy().sort_values("MonthNum") if not df.empty else pd.DataFrame()
df_wb_year_default = df_wb[df_wb["Year"] == year_default].copy().sort_values("MonthNum") if not df_wb.empty else pd.DataFrame()
df_dt_year_default = df_dt[df_dt["Year"] == year_default].copy().sort_values("MonthNum") if not df_dt.empty else pd.DataFrame()

latest_month_dt = None
if not df_year_default.empty:
    latest_month_dt = pd.to_datetime(df_year_default["Month"].max())
elif not df_wb_year_default.empty:
    latest_month_dt = pd.to_datetime(df_wb_year_default["Month"].max())
elif not df_dt_year_default.empty:
    latest_month_dt = pd.to_datetime(df_dt_year_default["Month"].max())
latest_label = latest_month_dt.strftime("%B %Y") if pd.notna(latest_month_dt) else str(year_default)

# =========================
# Header (logo + title + UPDATED menu)
# =========================
# =========================
# =========================
# Header (single gradient bar) + Tabs (directly below)
# =========================
st.markdown(
    f"""
    <style>
      /* One simple gradient header bar */
      .ppk-header--onebar {{
        background:
          linear-gradient(90deg, rgba(245,96,35,.18) 0%, rgba(245,96,35,.10) 42%, rgba(245,96,35,0) 78%),
          linear-gradient(180deg, {SLATE_900}, {SLATE_830});
        border: 1px solid rgba(255,255,255,.08);
        border-radius: 16px;
        padding: 12px 16px;
        box-shadow: 0 8px 24px rgba(0,0,0,.35);
      }}
      .ppk-title--big {{
        font-size: 1.8rem;
        font-weight: 900;
        letter-spacing: .2px;
        color: {TEXT};
        line-height: 1.15;
        margin: 2px 0 2px 0;
      }}
      /* Make the nav buttons tidy; these target only the immediate area below header */
      .ppk-tabs .stButton>button {{
        width: 100%;
        border: 1px solid rgba(255,255,255,.14);
        background: {SLATE_850};
        color: {TEXT};
        font-weight: 700;
        border-radius: 10px;
        padding: 8px 10px;
      }}
      .ppk-tabs .stButton>button:hover {{
        border-color: rgba(245,96,35,.55);
      }}
    </style>
    """,
    unsafe_allow_html=True
)

# --- Simple, self-contained header (no external CSS needed) ---
import base64

# Try to read the logo; safe fallback to None
_ppk_logo64 = None
try:
    with open("SARKS-02.png", "rb") as _f:
        _ppk_logo64 = base64.b64encode(_f.read()).decode("utf-8")
except Exception:
    _ppk_logo64 = None

with st.container():
    import base64
    from pathlib import Path

    # Point to your uploaded logo file
    logo_path = Path("ChatGPT Image Sep 2, 2025, 12_39_59 PM.png")  # make sure this file is in same folder as App_updated_SARKS.py
    logo_html = ""
    if logo_path.exists():
        with open(logo_path, "rb") as f:
            logo_b64 = base64.b64encode(f.read()).decode("utf-8")
        logo_html = f"<img src='data:image/png;base64,{logo_b64}' style='height:2.2em; width:auto; vertical-align:middle;'/>"

    header_html = f"""
    <div style="
        background:
          linear-gradient(90deg, rgba(245,96,35,.18) 0%, rgba(245,96,35,.10) 42%, rgba(245,96,35,0) 78%),
          linear-gradient(180deg, {SLATE_900}, {SLATE_830});
        border: 1px solid rgba(255,255,255,.08);
        border-radius: 16px;
        padding: 12px 16px;
        box-shadow: 0 8px 24px rgba(0,0,0,.35);
        ">
      <div style="
          display:flex; align-items:center; gap:10px;
          color:{TEXT}; font-size:1.8rem; font-weight:900; letter-spacing:.2px; line-height:1.2;
          white-space:nowrap; overflow:hidden; text-overflow:ellipsis;">
        {logo_html}
        <span style="color:{TEXT};">SARKS · IntelliLine · BOPET FilmLine KPI Dashboard</span>
      </div>
    </div>
    """

    st.markdown(header_html, unsafe_allow_html=True)


    # Tabs shown directly below the header bar (no separate subheader bar)
    st.markdown('<div class="ppk-tabs">', unsafe_allow_html=True)

    nav_labels = ["Home", "Monthly Update", "Forecast", "Waste Analytics", "Breakup", "KPIs Comparison", "Reports"]
    cols_nav = st.columns(len(nav_labels), gap="small")
    for i, label in enumerate(nav_labels):
        is_active = (st.session_state.active_page == label)
        btn_label = f"● {label}" if is_active else label
        if cols_nav[i].button(btn_label, key=f"nav_{label}", use_container_width=True):
            st.session_state.active_page = label

    st.markdown("</div>", unsafe_allow_html=True)

    # ---------- BLANK STATE (show when nothing is uploaded yet) ----------
if df.empty and df_wb.empty and df_dt.empty:
    st.markdown(
        """
        <div style="
            margin-top: 1rem;
            padding: 18px 20px;
            border: 1px solid rgba(255,255,255,.10);
            border-radius: 16px;
            background: linear-gradient(180deg, rgba(15,23,42,.55), rgba(15,23,42,.35));
            box-shadow: 0 8px 24px rgba(0,0,0,.35);
        ">
          <div style="font-size:1.1rem;font-weight:700;margin-bottom:.5rem;">
            👋 Upload your Excel workbook in the sidebar to begin
          </div>
          <div style="opacity:.9">
            <p style="margin:.2rem 0 .6rem 0;">XLSX preferred (XLS/CSV also ok). Use separate sheets:</p>
            <ul style="margin:.2rem 0 .6rem 1.2rem; line-height:1.5;">
              <li><b>KPI / Summary</b> — main monthly KPIs.</li>
              <li><b>Waste Breakup</b> — monthly waste by factors (Month column required).</li>
              <li><b>Downtime</b> — monthly downtime by factors (Month column required).</li>
            </ul>
            <p style="margin:.2rem 0 0 0; opacity:.85;">
              Tip: If your Month values don’t include the year (e.g. “Jan”, “Feb”), the app will infer the year.
            </p>
          </div>
        </div>
        """,
        unsafe_allow_html=True,
    )
    st.stop()
# ---------- END BLANK STATE ----------

# =========================
# HOME: YTD + MONTHLY tiles + charts
# =========================
def section_home(df_all: pd.DataFrame, df_year_fallback: pd.DataFrame) -> Dict[str, go.Figure]:
    charts: Dict[str, go.Figure] = {}
    if df_all.empty:
        st.info("Upload KPI file to view Home.")
        return charts

    kpi_years = sorted(df_all["Year"].dropna().astype(int).unique().tolist())
    y1, y2 = st.columns([1, 3])
    with y1:
        ytd_year = st.selectbox("Year for YTD/Key KPIs (Home)", kpi_years,
                                index=(kpi_years.index(year_default) if year_default in kpi_years else len(kpi_years)-1))
    with y2:
        if len(kpi_years) == 1:
            st.caption("Only one year found in the file — no alternate years to select.")

    df_year_local = df_all[df_all["Year"] == ytd_year].copy().sort_values("MonthNum")

    st.markdown(f"#### Year-to-Date KPIs — {ytd_year}")
    ycols = st.columns(6, gap="small")

    def ycard(slot, label, col, pct=False):
        val = yearly_value(df_year_local, col)
        agg = "SUM" if col in SUM_KPIS else ("AVG" if col in AVG_KPIS else "SUM")
        with slot:
            st.markdown(
                f"""
                <div class="year-card">
                  <div class="metric-caption">{label}</div>
                  <div class="year-val">{fmt_num(val, pct)}</div>
                  <span class="year-chip">{agg}</span>
                </div>
                """,
                unsafe_allow_html=True,
            )

    ycard(ycols[0], "Production (Tons)", "Production (Tons)")
    ycard(ycols[1], "Number of Accidents", "Number of Accidents")
    ycard(ycols[2], "Waste %", "Waste %", pct=True)
    ycard(ycols[3], "Total Down Time (Hours)", "Total Down Time (Hours)")
    ycard(ycols[4], "OEE %", "OEE %", pct=True)
    ycard(ycols[5], "Reclaim (Tons)", "Reclaim (Tons)")

    st.markdown("---")

    months_df = df_year_local[["MonthNum", "MonthName"]].dropna().drop_duplicates().sort_values("MonthNum")
    month_map = dict(zip(months_df["MonthName"], months_df["MonthNum"]))
    default_monthnum = int(months_df["MonthNum"].max()) if not months_df.empty else None
    m1, _m2 = st.columns([2, 2])
    with m1:
        sel_month_name = st.selectbox(
            f"Month for Key KPIs ({ytd_year})",
            options=list(month_map.keys()) if month_map else [],
            index=(list(month_map.values()).index(default_monthnum) if month_map and default_monthnum in month_map.values() else 0)
        ) if month_map else ""
    sel_monthnum = month_map.get(sel_month_name) if month_map else None

    sel_label = f"{sel_month_name} {ytd_year}" if sel_monthnum is not None else latest_label
    st.markdown(f"#### Key KPIs — {sel_label}")
    kpi_cols = st.columns(6, gap="small")

    def tile_at_month(slot, label, col, pct=False, inverse=False):
        if sel_monthnum is None:
            with slot: st.info("No months available for this year.")
            return
        k = kpi_delta_at_month(df_year_local, col, sel_monthnum)
        value = fmt_num(k["value"], pct)
        dv = k["delta"]
        delta_txt = "—" if (dv is None or pd.isna(dv)) else (f"{0.00:.2f}{'%' if pct else ''}" if abs(float(dv)) < 1e-12 else f"{dv:+.2f}{'%' if pct else ''}")
        with slot:
            st.markdown('<div class="metric-card">', unsafe_allow_html=True)
            st.markdown(f'<div class="metric-caption">{label}</div>', unsafe_allow_html=True)
            if delta_txt.startswith("0.00"):
                st.markdown(f'<div class="big-value">{value}</div>', unsafe_allow_html=True)
                st.markdown(f'<span class="delta-badge">= {delta_txt}</span>', unsafe_allow_html=True)
            else:
                st.metric(label=" ", value=value, delta=delta_txt, delta_color=("inverse" if inverse else "normal"))
            st.markdown('</div>', unsafe_allow_html=True)

    tile_at_month(kpi_cols[0], "Production (Tons)", "Production (Tons)")
    tile_at_month(kpi_cols[1], "Number of Accidents", "Number of Accidents")
    tile_at_month(kpi_cols[2], "Waste %", "Waste %", pct=True, inverse=True)
    tile_at_month(kpi_cols[3], "Total Down Time (Hours)", "Total Down Time (Hours)", inverse=True)
    tile_at_month(kpi_cols[4], "OEE %", "OEE %", pct=True)
    tile_at_month(kpi_cols[5], "Reclaim (Tons)", "Reclaim (Tons)")

    st.markdown("---")

    r1c1, r1c2 = st.columns(2, gap="large")
    with r1c1:
        if "Production (Tons)" in df_year_local.columns:
            st.markdown(f'<div class="ppk-panel"><div class="ppk-titlebar">Production — Monthly</div><div class="ppk-accent"></div>', unsafe_allow_html=True)
            charts["prod_line"] = line_fig(df_year_local, "Production (Tons)", "", COLORWAY[0], pct=False)
            st.plotly_chart(charts["prod_line"], use_container_width=True)
            st.markdown("</div>", unsafe_allow_html=True)
    with r1c2:
        if "Number of Accidents" in df_year_local.columns:
            st.markdown(f'<div class="ppk-panel"><div class="ppk-titlebar">Number of Accidents — Monthly</div><div class="ppk-accent"></div>', unsafe_allow_html=True)
            charts["accidents_bar"] = bar_fig(df_year_local, "Number of Accidents", "", COLORWAY[5])
            st.plotly_chart(charts["accidents_bar"], use_container_width=True)
            st.markdown("</div>", unsafe_allow_html=True)

    r2c1, r2c2 = st.columns(2, gap="large")
    with r2c1:
        donut_df, donut_err, _ = waste_composition_fixed(df_year_local)
        if donut_df is not None:
            st.markdown(f'<div class="ppk-panel"><div class="ppk-titlebar">Good Production vs Waste (YTD)</div><div class="ppk-accent"></div>', unsafe_allow_html=True)
            charts["donut"] = px.pie(
                donut_df, names="Category", values="Tons", hole=0.55,
                color="Category",
                color_discrete_map={"Good Production": "#34D399", "Waste": "#EF4444"}
            )
            charts["donut"].update_traces(textinfo="percent+label", pull=[0, 0.06])
            st.plotly_chart(_apply_theme(charts["donut"]), use_container_width=True)
            st.caption(f"Total Production (context) = Good Production + Waste = {fmt_num(donut_df['Tons'].sum())} tons")
            st.markdown("</div>", unsafe_allow_html=True)
        else:
            st.info(donut_err or "Insufficient data for donut chart.")
    with r2c2:
        if {"Production (Tons)", "Waste (Tons)"}.issubset(df_year_local.columns):
            st.markdown(f'<div class="ppk-panel"><div class="ppk-titlebar">Total Production by Month (Stacked)</div><div class="ppk-accent"></div>', unsafe_allow_html=True)
            charts["stacked"] = stacked_bar_fig(df_year_local, "Production (Tons)", "Waste (Tons)", "", "#60A5FA", "#F59E0B")
            st.plotly_chart(charts["stacked"], use_container_width=True)
            st.markdown("</div>", unsafe_allow_html=True)

    r3c1, r3c2 = st.columns(2, gap="large")
    with r3c1:
        if "Total Down Time (Hours)" in df_year_local.columns:
            st.markdown(f'<div class="ppk-panel"><div class="ppk-titlebar">Downtime — Monthly</div><div class="ppk-accent"></div>', unsafe_allow_html=True)
            charts["downtime_area"] = area_fig(df_year_local, "Total Down Time (Hours)", "", "#06B6D4")
            st.plotly_chart(charts["downtime_area"], use_container_width=True)
            st.markdown("</div>", unsafe_allow_html=True)
    with r3c2:
        if "OEE %" in df_year_local.columns:
            st.markdown(f'<div class="ppk-panel"><div class="ppk-titlebar">OEE% — Monthly</div><div class="ppk-accent"></div>', unsafe_allow_html=True)
            charts["oee_area"] = area_fig(df_year_local, "OEE %", "", COLORWAY[3], pct=True)
            st.plotly_chart(charts["oee_area"], use_container_width=True)
            st.markdown("</div>", unsafe_allow_html=True)

    r4c1, r4c2 = st.columns([1.05, 0.95], gap="large")
    with r4c1:
        if "Reclaim (Tons)" in df_year_local.columns:
            st.markdown(f'<div class="ppk-panel"><div class="ppk-titlebar">Reclaim — Monthly</div><div class="ppk-accent"></div>', unsafe_allow_html=True)
            charts["reclaim_lollipop"] = lollipop_fig(df_year_local, "Reclaim (Tons)", "", TEAL)
            st.plotly_chart(charts["reclaim_lollipop"], use_container_width=True)
            st.markdown("</div>", unsafe_allow_html=True)
    with r4c2:
        if "Q2 (Tons)" in df_year_local.columns:
            st.markdown(f'<div class="ppk-panel"><div class="ppk-titlebar">Quality — Q2 Tons (Monthly)</div><div class="ppk-accent"></div>', unsafe_allow_html=True)
            charts["q2_bar"] = bar_fig(df_year_local, "Q2 (Tons)", "", COLORWAY[8])
            st.plotly_chart(charts["q2_bar"], use_container_width=True)
            st.markdown("</div>", unsafe_allow_html=True)

    st.markdown(f'<div class="ppk-panel"><div class="ppk-titlebar">KPI Correlation</div><div class="ppk-accent" style="width:140px;"></div>', unsafe_allow_html=True)
    heat = corr_heatmap(df_year_local, "")
    if heat is not None:
        charts["corr"] = heat
        st.plotly_chart(heat, use_container_width=True)
    else:
        st.caption("Not enough numeric columns to compute correlation.")
    st.markdown("</div>", unsafe_allow_html=True)

    # Downloads (Home)
    figs = [charts[k] for k in ["prod_line","accidents_bar","donut","stacked","downtime_area","oee_area","reclaim_lollipop","q2_bar","corr"] if k in charts]
    narrative = [
        f"Year-to-date view for {ytd_year}.",
        "Tiles show YTD totals/averages and last-month deltas.",
        "Charts highlight trends across production, waste, downtime, OEE, and quality."
    ]
    _download_row("Home", df_year_local, figs, narrative)

    return charts

# =========================
# QUALITY / WASTE / REPORTS (unchanged core + downloads)
# =========================
def section_quality(df_year: pd.DataFrame) -> Dict[str, go.Figure]:
    charts: Dict[str, go.Figure] = {}
    if df_year.empty:
        st.info("Upload KPI file to view Quality.")
        return charts
    if "Q2 (Tons)" in df_year.columns:
        charts["q2_bar"] = bar_fig(df_year, "Q2 (Tons)", "Q2 — Monthly", COLORWAY[8])
        st.plotly_chart(charts["q2_bar"], use_container_width=True)
    else:
        st.info("No dedicated quality-specific KPI columns detected (e.g., 'Q2 (Tons)').")
    return charts


def section_waste(df_year: pd.DataFrame) -> Dict[str, go.Figure]:
    charts: Dict[str, go.Figure] = {}
    if df_year.empty:
        st.info("Upload KPI file to view Waste Analytics.")
        return charts

    if {"Production (Tons)", "Waste (Tons)"}.issubset(df_year.columns):
        c1, c2 = st.columns([1.1, 1.3], gap="large")

        # Left: Good Production vs Waste (YTD) donut
        with c1:
            donut_df, donut_err, _ = waste_composition_fixed(df_year)
            if donut_df is not None:
                charts["donut"] = px.pie(
                    donut_df,
                    names="Category",
                    values="Tons",
                    hole=0.58,
                    color="Category",
                    color_discrete_map={"Good Production": "#34D399", "Waste": "#EF4444"},
                )
                charts["donut"].update_traces(textinfo="percent+label", pull=[0, 0.08])
                st.plotly_chart(_apply_theme(charts["donut"]), use_container_width=True)
            else:
                st.info(donut_err or "Insufficient data for donut chart.")

        # Right: Total Production by Month (Stacked)
        with c2:
            charts["stacked"] = stacked_bar_fig(
                df_year,
                "Production (Tons)",
                "Waste (Tons)",
                "Total Production by Month (Stacked)",
                "#60A5FA",
                "#F59E0B",
            )
            st.plotly_chart(charts["stacked"], use_container_width=True)
    else:
        st.info("Upload data with 'Production (Tons)' and 'Waste (Tons)' to view waste analytics.")

    return charts

# BREAKUP (Waste & Downtime)
# =========================
# ----- Waste helpers (UPDATED) -----
from typing import Optional

def _wb_coerce_month(dt_like: pd.Series, year: Optional[int]) -> tuple[pd.Series, pd.Series, pd.Series, pd.Series]:
    """
    Make Month/Year/MonthNum/MonthName from a 'Month' column for WASTE sheets.
    - If text includes a year (e.g., 'Jan 2024' or '2024-01'), use that year.
    - Else use `year` if provided (assumed_year).
    - Else fall back safely to the current year.
    """
    import pandas as pd

    # If it's already datetime, keep its year & month
    if pd.api.types.is_datetime64_any_dtype(dt_like):
        Month = pd.to_datetime(dt_like, errors="coerce")
        Year = Month.dt.year.astype("Int64")
        MonthNum = Month.dt.month.astype(int)
        MonthName = Month.dt.strftime("%b")
        return Month, Year, MonthNum, MonthName

    s = dt_like.astype(str).str.strip()

    # Month name → number maps
    map_short = {m: i for i, m in enumerate(
        ["Jan","Feb","Mar","Apr","May","Jun","Jul","Aug","Sep","Oct","Nov","Dec"], start=1)}
    mon_num = s.str[:3].str.title().map(map_short)

    full_map = {m: i for i, m in enumerate(
        ["January","February","March","April","May","June","July","August",
         "September","October","November","December"], start=1)}
    mon_num = mon_num.fillna(s.str.title().map(full_map))

    # Try to parse full dates to grab year/month if present
    dt = pd.to_datetime(s, errors="coerce", dayfirst=False, yearfirst=False)

    # Prefer parsed year; then any 4-digit year in text; then assumed year
    yy = dt.dt.year.astype("Float64")  # allow NaN
    try:
        yr_txt = pd.to_numeric(s.str.extract(r"(\d{4})")[0], errors="coerce")
        yy = yy.fillna(yr_txt)
    except Exception:
        pass
    if year is not None:
        yy = yy.fillna(float(year))

    # Final fallback if still missing everywhere
    yy = yy.ffill().bfill()
    if yy.isna().all():
        yy = pd.Series([pd.Timestamp.today().year] * len(s), index=s.index, dtype="Float64")

    # Month numbers: from our maps or parsed dt; default to 1
    mm = mon_num.fillna(dt.dt.month).fillna(1).astype(int)

    Year = yy.astype("Int64")
    Month = pd.to_datetime(
        {"year": Year.fillna(pd.Timestamp.today().year).astype(int), "month": mm, "day": 1},
        errors="coerce"
    )
    MonthName = Month.dt.strftime("%b")
    return Month, Year, mm, MonthName


def _wb_numeric(s: pd.Series) -> pd.Series:
    s2 = s.astype(str).str.replace(",", "", regex=False).str.replace("%", "", regex=False)
    return pd.to_numeric(s2, errors="coerce")


def _wb_clean(df: pd.DataFrame, assumed_year: Optional[int]):
    """
    Clean 'Waste Breakup' style data.
    - Respects an existing Year column (doesn't overwrite).
    - Infers year from text (e.g. 'Jan 2024') and only uses `assumed_year` where year is missing.
    - Returns: (clean_df, total_col_name, factor_columns_list)
    """
    if df.empty:
        return df, None, []

    # If first row is headers (Unnamed...), elevate it
    if any(str(c).startswith("Unnamed") for c in df.columns):
        new_cols = list(df.iloc[0].values)
        df = df.iloc[1:].copy()
        df.columns = new_cols

    # Standardize
    df.columns = [str(c).strip() for c in df.columns]
    if "Month" not in df.columns:
        df.rename(columns={df.columns[0]: "Month"}, inplace=True)

    # Keep a copy of any pre-existing Year to prefer it later
    orig_year_col = None
    if "Year" in df.columns:
        orig_year_col = pd.to_numeric(df["Year"], errors="coerce")

    # Parse Month/Year
    Month, Year, MonthNum, MonthName = _wb_coerce_month(df["Month"], assumed_year)

    # If the sheet already had a Year column with values, prefer it
    if orig_year_col is not None:
        Year = orig_year_col.fillna(Year)

    # Rebuild Month from the final Year + MonthNum for consistency
    Month = pd.to_datetime(
        {"year": Year.fillna(pd.Timestamp.today().year).astype(int), "month": MonthNum, "day": 1},
        errors="coerce"
    )
    MonthName = Month.dt.strftime("%b")

    # Assign back
    df["Month"] = Month
    df["Year"] = Year
    df["MonthNum"] = MonthNum
    df["MonthName"] = MonthName

    # Make numerics (skip date fields)
    for c in df.columns:
        if c in {"Month", "Year", "MonthNum", "MonthName"}:
            continue
        df[c] = _wb_numeric(df[c])

    # Column ordering
    present = [c for c in WB_EXPECTED_COLUMNS if c in df.columns and c not in {"Month","Year","MonthNum","MonthName"}]
    others = [c for c in df.columns if c not in present and c not in {"Month","Year","MonthNum","MonthName"}]
    order = ["Month","Year","MonthNum","MonthName"] + present + others
    order = list(dict.fromkeys([c for c in order if c in df.columns]))
    df = df[order].sort_values("MonthNum").reset_index(drop=True)

    # Total column detection + factor list
    total_col = "Monthly Waste (Tons)" if "Monthly Waste (Tons)" in df.columns else detect_total_waste_col(df.columns)
    factor_order = ["Process","Planning(G.C/M.C)","Special Films","Electrical","Mechanical","Power House","No Job Order"]
    factors = [c for c in factor_order if c in df.columns]
    extras = [
        c for c in df.columns
        if (c not in {"Month","Year","MonthNum","MonthName","Waste %", total_col})
        and pd.api.types.is_numeric_dtype(df[c])
        and c not in factors
    ]
    factors = factors + extras
    return df, total_col, factors


# ----- Downtime helpers (UNCHANGED) -----
def _dt_coerce_month(dt_like: pd.Series, year: int):
    if pd.api.types.is_datetime64_any_dtype(dt_like):
        Month = pd.to_datetime(dt_like)
        MonthName = Month.dt.strftime("%b")
        return Month, Month.dt.year, Month.dt.month, MonthName
    s = dt_like.astype(str).str.strip()
    map_short = {m:i for i,m in enumerate(["Jan","Feb","Mar","Apr","May","Jun","Jul","Aug","Sep","Oct","Nov","Dec"], start=1)}
    mon_num = s.str[:3].str.title().map(map_short)
    full_map = {m:i for i,m in enumerate(
        ["January","February","March","April","May","June","July","August","September","October","November","December"], start=1)}
    mon_num = mon_num.fillna(s.str.strip().str.title().map(full_map))
    dt = pd.to_datetime(s, errors="coerce")
    mon_num = mon_num.fillna(dt.dt.month)
    yy = dt.dt.year.fillna(year)
    mm = mon_num.fillna(1).astype(int)
    Month = pd.to_datetime(pd.DataFrame({"year": yy.fillna(year).astype(int), "month": mm, "day": 1}))
    MonthName = Month.dt.strftime("%b")
    return Month, Month.dt.year, Month.dt.month, MonthName

def _dt_numeric(s: pd.Series) -> pd.Series:
    s2 = s.astype(str).str.replace(",","", regex=False).str.replace("%","", regex=False)
    return pd.to_numeric(s2, errors="coerce")

def _dt_clean(df: pd.DataFrame, assumed_year: int):
    if df.empty:
        return df, None, []
    if any(str(c).startswith("Unnamed") for c in df.columns):
        new_cols = list(df.iloc[0].values)
        df = df.iloc[1:].copy()
        df.columns = new_cols

    df.columns = [str(c).strip() for c in df.columns]
    if "Month" not in df.columns:
        df.rename(columns={df.columns[0]: "Month"}, inplace=True)

    Month, Year, MonthNum, MonthName = _dt_coerce_month(df["Month"], assumed_year)
    df["Month"] = Month; df["Year"] = Year; df["MonthNum"] = MonthNum; df["MonthName"] = MonthName

    for c in df.columns:
        if c in {"Month","Year","MonthNum","MonthName"}:
            continue
        df[c] = _dt_numeric(df[c])

    present = [c for c in DT_EXPECTED_COLUMNS if c in df.columns and c not in {"Month","Year","MonthNum","MonthName"}]
    others = [c for c in df.columns if c not in present and c not in {"Month","Year","MonthNum","MonthName"}]
    order = ["Month","Year","MonthNum","MonthName"] + present + others
    order = list(dict.fromkeys([c for c in order if c in df.columns]))
    df = df[order].sort_values("MonthNum").reset_index(drop=True)

    total_col = "Total Down Time (Hours)" if "Total Down Time (Hours)" in df.columns else None
    factor_order = [
        "No Job order",
        "Planned Preventive Maintenance+Filter Change",
        "Process",
        "Planning (G.C/M.C)",
        "Mech. Plant",
        "Mech. Utilities",
        "Electrical Plant",
        "Electrical Utilities",
        "Power House4",
        "Rental Power",
    ]
    factors = [c for c in factor_order if c in df.columns]
    extras = [c for c in df.columns if (c not in {"Month","Year","MonthNum","MonthName", total_col}) and pd.api.types.is_numeric_dtype(df[c]) and c not in factors]
    factors = factors + extras
    return df, total_col, factors

def section_breakup_combined():
    """
    Breakup tab with per-page Year/Month filtering that drives ALL visuals.
    - Waste: stacked factors, shares, Pareto & donut, correlation
    - Downtime: same structure (Hours)
    NOTE:
    - Nothing renders before "Breakup — Choose View".
    - Waste uses the same Year/Month selection behavior as Downtime.
    - Unique namespaces prevent widget/state collisions across the app.
    """

    import numpy as np
    import pandas as pd
    import plotly.express as px
    import plotly.graph_objects as go
    import streamlit as st

    # -------------------------
    # 0) Header & mode toggle (nothing above this line should render)
    # -------------------------
    if "breakup_mode" not in st.session_state:
        st.session_state.breakup_mode = "Waste"

    st.markdown("#### Breakup — Choose View")
    sub_cols = st.columns(2, gap="small")
    for i, label in enumerate(["Waste", "Downtime"]):
        is_on = (st.session_state.breakup_mode == label)
        btn = f"● {label}" if is_on else label
        if sub_cols[i].button(btn, key=f"breakup_tab_{label}", use_container_width=True):
            st.session_state.breakup_mode = label

    mode = st.session_state.breakup_mode

    # -------------------------
    # 1) Reset stale state on tab switch (only our page keys)
    # -------------------------
    if "_last_breakup_mode" not in st.session_state:
        st.session_state["_last_breakup_mode"] = mode
    if st.session_state["_last_breakup_mode"] != mode:
        for k in list(st.session_state.keys()):
            if k.startswith("br_"):   # our page/module namespace
                del st.session_state[k]
        st.session_state["_last_breakup_mode"] = mode

    # Unique namespace per page + mode (prevents collisions with other pages)
    active_page = st.session_state.get("active_page", "root").lower()
    page_ns = f"br_{active_page}"
    mode_ns = "waste" if mode == "Waste" else "dt"
    ns = f"{page_ns}_{mode_ns}"          # e.g., br_home_waste
    picker_ns = f"{ns}_picker"           # for pick_year_month widgets

    # -------------------------
    # 2) Load & clean data (per mode)
    # -------------------------
    if mode == "Waste":
        wb_raw = st.session_state.get("waste_raw_df", pd.DataFrame())
        if wb_raw.empty:
            st.info("👋 No Waste Breakup sheet detected. Add a 'Waste Breakup' sheet to your workbook.")
            return {}

        # IMPORTANT: let cleaner infer years from data (do NOT force a single year)
        df_all, total_col, factor_cols_all = _wb_clean(wb_raw.copy(), None)
        unit_lbl = "t"
        pct_label = "Normalize stacked to % of total (waste)"
        label_prefix = "Breakup (Waste) — "
    else:
        dt_raw = st.session_state.get("downtime_raw_df", pd.DataFrame())
        if dt_raw.empty:
            st.info("👋 No Downtime Breakup sheet detected. Add a 'Downtime Breakup' sheet to your workbook.")
            return {}

        # Downtime already works well; keep same behavior (infer years)
        df_all, total_col, factor_cols_all = _dt_clean(dt_raw.copy(), None)
        unit_lbl = "h"
        pct_label = "Normalize stacked to % of total (downtime)"
        label_prefix = "Breakup (Downtime) — "

    # Defensive dedupe + numeric types
    df_all = df_all.loc[:, ~df_all.columns.duplicated()]
    if "Year" in df_all.columns:
        df_all["Year"] = pd.to_numeric(df_all["Year"], errors="coerce").astype("Int64")
    if "MonthNum" in df_all.columns:
        df_all["MonthNum"] = pd.to_numeric(df_all["MonthNum"], errors="coerce")

    # -------------------------
    # 3) Single Year/Month control INSIDE the active view
    # -------------------------
    _, sel_monthnum, df_vis = pick_year_month(df_all, label_prefix, key_ns=picker_ns)
    if df_vis is None or df_vis.empty:
        st.info("No rows for the selected Year/Month.")
        return {}

    # Factors present in the visible slice (intersection only)
    factor_cols = [
        c for c in factor_cols_all
        if c in df_vis.columns and c not in ("Year", "MonthNum", "MonthName")
    ]

    # -------------------------
    # 4) Controls (namespaced keys so they never collide)
    # -------------------------
    opt1, opt2, opt3 = st.columns([1.2, 1.2, 1], gap="small")
    with opt1:
        normalize_pct = st.checkbox(pct_label, value=False, key=f"{ns}_pct")
    with opt2:
        top_n = st.slider(f"Pareto Top-N ({mode.lower()})", 3, 12, 6, key=f"{ns}_topn")
    with opt3:
        default_sel = factor_cols[: min(6, len(factor_cols))]
        chosen = st.multiselect("Select factors", options=factor_cols, default=default_sel, key=f"{ns}_chosen")

    if not chosen:
        st.info("Select at least one factor to visualize.")
        return {}

    # -------------------------
    # 5) Consistency check (Total vs Sum of chosen factors)
    # -------------------------
    df_vis["__SumFactors__"] = df_vis[chosen].sum(axis=1, skipna=True) if chosen else 0.0
    gap_text = ""
    if total_col and total_col in df_vis.columns:
        df_vis["__Diff__"] = (df_vis[total_col] - df_vis["__SumFactors__"]).fillna(0.0)
        max_abs_diff = float(df_vis["__Diff__"].abs().max())
        gap_text = f"Max Δ (Total − Sum of factors): {max_abs_diff:.2f} {unit_lbl}"
        if max_abs_diff > 1e-6:
            st.warning(f"Breakup check — {gap_text}")
    else:
        total_col = "__SumFactors__"
        st.info("No explicit total column found; using Sum of factors as total.")

    # -------------------------
    # 6) Metrics (based on current selection – fixes “always Dec”)
    # -------------------------
    latest = df_vis.sort_values("MonthNum").tail(1)
    latest_label = latest.iloc[0]["MonthName"] if not latest.empty else ""
    latest_total = float(latest[total_col].values[0]) if not latest.empty else np.nan
    ytd_total = float(pd.to_numeric(df_vis[total_col], errors="coerce").fillna(0).sum()) if total_col in df_vis.columns else np.nan

    st.markdown("#### Key Highlights")
    m1, m2, m3, m4 = st.columns(4, gap="small")
    with m1:
        st.markdown(
            f"""
            <div class="metric-card">
              <div class="metric-caption">Latest (within selection)</div>
              <div class="big-value">{latest_label or "—"}</div>
              <span class="delta-badge">{gap_text if gap_text else "Validated"}</span>
            </div>
            """,
            unsafe_allow_html=True,
        )
    with m2:
        st.markdown(
            f"""
            <div class="metric-card">
              <div class="metric-caption">Total {mode} (Latest)</div>
              <div class="big-value">{(latest_total if not np.isnan(latest_total) else 0):,.2f} {unit_lbl}</div>
              <span class="delta-badge">{total_col}</span>
            </div>
            """,
            unsafe_allow_html=True,
        )
    with m3:
        if chosen and not latest.empty:
            latest_f = latest[chosen].iloc[0].sort_values(ascending=False)
            top_name = str(latest_f.index[0]); top_val = float(latest_f.iloc[0])
        else:
            top_name, top_val = "—", 0.0
        st.markdown(
            f"""
            <div class="metric-card">
              <div class="metric-caption">Top Contributor (Latest)</div>
              <div class="big-value">{top_name}</div>
              <span class="delta-badge">{top_val:,.2f} {unit_lbl}</span>
            </div>
            """,
            unsafe_allow_html=True,
        )
    with m4:
        st.markdown(
            f"""
            <div class="metric-card">
              <div class="metric-caption">YTD {mode} (within selection)</div>
              <div class="big-value">{(ytd_total if not np.isnan(ytd_total) else 0):,.2f} {unit_lbl}</div>
              <span class="delta-badge">SUM</span>
            </div>
            """,
            unsafe_allow_html=True,
        )

    st.markdown("---")

    # -------------------------
    # 7) Charts — always built from df_vis (filtered slice)
    # -------------------------
    order = month_order_layout()
    figs_for_pdf = []

    # 7.1) Total trend
    if total_col in df_vis.columns:
        fig_total = px.line(
            df_vis.sort_values("MonthNum"),
            x="MonthName", y=total_col, markers=True,
            title=f"Total {mode} — Monthly Trend (Selection)"
        )
        fig_total.update_traces(line=dict(width=3))
        fig_total = _apply_theme(fig_total).update_layout(yaxis_title=("Tons" if mode == "Waste" else "Hours"))
        st.plotly_chart(fig_total, use_container_width=True)
        figs_for_pdf.append(fig_total)

    # 7.2) Stacked bars
    plot_df = df_vis[["MonthName", "MonthNum", total_col] + chosen].sort_values("MonthNum").copy()
    ylab = ("Tons" if mode == "Waste" else "Hours")
    if normalize_pct:
        for c in chosen:
            plot_df[c] = np.where(plot_df[total_col] > 0, (plot_df[c] / plot_df[total_col]) * 100.0, 0.0)
        ylab = f"% of Total {mode}"

    stacked = go.Figure()
    for c in chosen:
        stacked.add_bar(name=c, x=plot_df["MonthName"], y=plot_df[c])
    stacked.update_layout(
        title=f"{mode} Factors — Monthly Stacked (Selection)",
        barmode="stack", xaxis=order, xaxis_title="Month", yaxis_title=ylab,
        hovermode="x unified",
        legend=dict(orientation="h", yanchor="bottom", y=1.02),
    )
    stacked = _apply_theme(stacked)
    st.plotly_chart(stacked, use_container_width=True)
    figs_for_pdf.append(stacked)

    # 7.3) Shares area (optional)
    show_shares = st.checkbox("Show factor shares area (as % of total)", value=True, key=f"{ns}_shares")
    if show_shares:
        share_df = df_vis[["MonthName", "MonthNum", total_col] + chosen].sort_values("MonthNum").copy()
        for c in chosen:
            share_df[c] = np.where(share_df[total_col] > 0, (share_df[c] / share_df[total_col]) * 100.0, 0.0)
        melted = share_df.melt(id_vars=["MonthName", "MonthNum"], value_vars=chosen, var_name="Factor", value_name="SharePct")
        area = px.area(
            melted.sort_values("MonthNum"), x="MonthName", y="SharePct", color="Factor",
            title=f"Factor Shares of Total {mode} — Over Time (%)", groupnorm="percent"
        )
        area.update_layout(xaxis=order, yaxis_title=f"% of Total {mode}")
        area = _apply_theme(area)
        st.plotly_chart(area, use_container_width=True)
        figs_for_pdf.append(area)

    # 7.4) Pareto + 7.5) Donut — latest visible month
    val_col = ("Tons" if mode == "Waste" else "Hours")
    if not latest.empty:
        latest_vals = latest[chosen].iloc[0].sort_values(ascending=False)
        pareto_df = latest_vals.reset_index(); pareto_df.columns = ["Factor", val_col]
        pareto_df = pareto_df.head(top_n)
    else:
        pareto_df = pd.DataFrame(columns=["Factor", val_col])

    pareto = px.bar(pareto_df, x="Factor", y=val_col, title=f"Pareto — Top {top_n} Contributors ({latest_label})")
    pareto = _apply_theme(pareto)
    st.plotly_chart(pareto, use_container_width=True)
    figs_for_pdf.append(pareto)

    if not latest.empty:
        donut_df = latest[chosen].T.reset_index()
        donut_df.columns = ["Factor", val_col]
        donut = px.pie(donut_df, names="Factor", values=val_col, hole=0.55)
        donut.update_traces(textinfo="percent+label", pull=[0.04] * len(donut_df))
        donut.update_layout(title=f"Composition — {latest_label}")
        donut = _apply_theme(donut)
        st.plotly_chart(donut, use_container_width=True)
        figs_for_pdf.append(donut)

    # 7.6) Correlation
    num_cols = df_vis[chosen + ([total_col] if total_col else [])].select_dtypes(include=[np.number])
    if num_cols.shape[1] >= 2:
        corr = num_cols.corr()
        heat = go.Figure(data=go.Heatmap(
            z=corr.values, x=list(corr.columns), y=list(corr.index),
            zmin=-1, zmax=1, colorscale="Blues"
        ))
        heat.update_layout(title=f"Correlation: Selected Factors (incl. Total) — Selection Only")
        heat = _apply_theme(heat)
        st.plotly_chart(heat, use_container_width=True)
        figs_for_pdf.append(heat)

    # -------------------------
    # 8) Downloads
    # -------------------------
    _download_row(f"Breakup_{mode}", df_vis, figs_for_pdf, [
        f"{mode} breakup by factor across selected months.",
        "Pareto highlights dominant contributors for latest visible month.",
        "Composition donut and correlation reflect the selection."
    ])

    return {}

# =========================
# Reports (overview + downloads)
# =========================
def section_reports(df_all: pd.DataFrame, df_wb_all: pd.DataFrame, df_dt_all: pd.DataFrame):
    st.markdown("#### Report — Uploaded Data Overview")

    figs = []
    if not df_all.empty and "Production (Tons)" in df_all.columns:
        dfp = df_all.copy().sort_values(["Year","MonthNum"])
        pfig = px.line(dfp, x="Month", y="Production (Tons)", title="Production Over Time", markers=True)
        figs.append(_apply_theme(pfig))

    if not df_all.empty:
        with st.expander("KPI Data — All Rows (cleaned)", expanded=True):
            st.dataframe(df_all, use_container_width=True, hide_index=True)
            st.download_button("⬇ CSV — KPI (all rows)", data=_csv_bytes(df_all),
                               file_name=f"kpi_all_{datetime.now():%Y%m%d}.csv", use_container_width=True)
    else:
        st.info("No KPI data detected.")

    if not df_wb_all.empty:
        with st.expander("Waste Breakup — All Rows (cleaned)", expanded=False):
            st.dataframe(df_wb_all, use_container_width=True, hide_index=True)
            st.download_button("⬇ CSV — Waste Breakup", data=_csv_bytes(df_wb_all),
                               file_name=f"waste_breakup_all_{datetime.now():%Y%m%d}.csv", use_container_width=True)
    else:
        st.info("No Waste Breakup data detected.")

    if not df_dt_all.empty:
        with st.expander("Downtime Breakup — All Rows (cleaned)", expanded=False):
            st.dataframe(df_dt_all, use_container_width=True, hide_index=True)
            st.download_button("⬇ CSV — Downtime Breakup", data=_csv_bytes(df_dt_all),
                               file_name=f"downtime_breakup_all_{datetime.now():%Y%m%d}.csv", use_container_width=True)
    else:
        st.info("No Downtime Breakup data detected.")

    _download_row("Reports", df_all, figs, [
        "This report compiles all uploaded and cleaned tables.",
        "Use CSVs above for data extraction; PDF/PPT include overview chart(s)."
    ])
    return {}

# =========================
# KPIs Comparison (with downloads)
# =========================
BASE_KPIS = [
    "Production (Tons)", "Number of Accidents", "Waste (Tons)",
    "Total Down Time (Hours)", "OEE %", "Reclaim (Tons)", "Waste %"
]
def section_kpi_comparison(df_year: pd.DataFrame):
    if df_year.empty:
        st.info("Upload KPI file to view KPIs Comparison.")
        return {}
    charts = {}
    figs = []
    avail = [c for c in BASE_KPIS if c in df_year.columns]
    for base in avail:
        others = [c for c in avail if c != base]
        if not others:
            continue
        st.markdown(f"##### {base} — comparison")
        sel_multi = st.multiselect(f"Compare {base} with", options=others, default=[others[0]] if others else [], key=f"cmp_{base}")
        if not sel_multi:
            st.info("Pick at least one KPI to compare.")
            continue
        fig = compare_multi_kpis_fig(df_year, base, sel_multi)
        if fig is not None:
            charts[f"cmp_{base}"] = fig
            st.plotly_chart(fig, use_container_width=True)
            figs.append(fig)

    _download_row("KPIs_Comparison", df_year, figs, [
        "Each chart compares a base KPI against other KPIs.",
        "Percent KPIs are placed on a secondary axis when useful."
    ])
    return charts

# =========================
# Monthly Update (with downloads)
# =========================
def section_monthly_update(df_all: pd.DataFrame):
    if df_all.empty:
        st.info("Upload KPI file to view Monthly Update.")
        return {}
    charts = {}
    years = sorted(df_all["Year"].dropna().astype(int).unique().tolist())
    c1, c2 = st.columns(2)
    with c1:
        ysel = st.selectbox("Year", years, index=(years.index(year_default) if year_default in years else len(years)-1))
    months_df = df_all[df_all["Year"] == ysel][["MonthNum","MonthName"]].drop_duplicates().sort_values("MonthNum")
    month_map = dict(zip(months_df["MonthName"], months_df["MonthNum"]))
    with c2:
        msel_name = st.selectbox("Month", list(month_map.keys()) if month_map else [])
    if not month_map:
        st.info("No months available.")
        return charts
    msel = month_map[msel_name]
    row = df_all[(df_all["Year"] == ysel) & (df_all["MonthNum"] == msel)]
    if row.empty:
        st.info("No data for the selected month.")
        return charts

    figs = []
    abs_cols = [c for c in ["Production (Tons)","Waste (Tons)","Total Down Time (Hours)","Reclaim (Tons)","Number of Accidents"] if c in row.columns]
    pct_cols = [c for c in ["OEE %","Waste %"] if c in row.columns]

    if abs_cols:
        vals = row.iloc[0][abs_cols].astype(float)
        df_abs = pd.DataFrame({"KPI": abs_cols, "Value": vals.values})
        st.markdown(f'<div class="ppk-panel"><div class="ppk-titlebar">Monthly Detail — {msel_name} {ysel} (Absolutes)</div><div class="ppk-accent"></div>', unsafe_allow_html=True)
        charts["month_abs"] = px.bar(df_abs, x="KPI", y="Value", title="", color="KPI",
                                     color_discrete_sequence=COLORWAY[:max(3, len(abs_cols))])
        st.plotly_chart(_apply_theme(charts["month_abs"]), use_container_width=True)
        st.markdown("</div>", unsafe_allow_html=True)
        figs.append(charts["month_abs"])

    if pct_cols:
        vals = row.iloc[0][pct_cols].astype(float)
        df_pct = pd.DataFrame({"KPI": pct_cols, "Value": vals.values})
        st.markdown(f'<div class="ppk-panel"><div class="ppk-titlebar">Monthly Detail — {msel_name} {ysel} (% KPIs)</div><div class="ppk-accent"></div>', unsafe_allow_html=True)
        charts["month_pct"] = px.bar(df_pct, x="KPI", y="Value", title="", color="KPI",
                                     color_discrete_sequence=COLORWAY[:max(3, len(pct_cols))])
        charts["month_pct"].update_yaxes(ticksuffix=" %")
        st.plotly_chart(_apply_theme(charts["month_pct"]), use_container_width=True)
        st.markdown("</div>", unsafe_allow_html=True)
        figs.append(charts["month_pct"])

    mtable = pd.concat([df_abs if abs_cols else pd.DataFrame(), df_pct if pct_cols else pd.DataFrame()], ignore_index=True) if (abs_cols or pct_cols) else pd.DataFrame()
    _download_row("Monthly_Update", mtable, figs, [
        f"Snapshot of {msel_name} {ysel} KPIs in absolute and percentage terms."
    ])
    return charts

# =========================
# Forecast (Production RANGE + residual spread bands)
# =========================

def _linfit_y_on_prod(df_all: pd.DataFrame, y_col: str) -> Optional[Tuple[float, float, float]]:
    """
    Return (slope a, intercept b, residual_std) for y ≈ a*Prod + b using all rows with both values present.
    Uses 'Production (Tons)' as the single driver.
    """
    if df_all.empty or "Production (Tons)" not in df_all.columns or y_col not in df_all.columns:
        return None
    tmp = df_all[["Production (Tons)", y_col]].dropna()
    if len(tmp) < 2:
        return None
    x = tmp["Production (Tons)"].astype(float).values
    y = tmp[y_col].astype(float).values
    try:
        a, b = np.polyfit(x, y, 1)
        resid = y - (a * x + b)
        std = float(pd.Series(resid).std(ddof=1)) if len(resid) > 1 else 0.0
        return float(a), float(b), std
    except Exception:
        return None

def _clip_pct(v: float) -> float:
    return max(0.0, min(100.0, float(v)))

def _predict_range_from_prod(
    df_all: pd.DataFrame, prod_lo: float, prod_hi: float
) -> Dict[str, Dict[str, Optional[float]]]:
    """
    For each KPI, predict a midpoint and a (lo, hi) band from a production range.
    Band combines the effect of the production range + ~80% residual spread (z≈1.28).
    """
    out: Dict[str, Dict[str, Optional[float]]] = {}
    KPIS = ["Waste (Tons)", "OEE %", "Total Down Time (Hours)", "Reclaim (Tons)", "Waste %", "Number of Accidents"]
    mid_prod = (prod_lo + prod_hi) / 2.0

    for col in KPIS:
        out[col] = {"mid": None, "lo": None, "hi": None}
        # Skip Waste % direct fit; derive from predicted Waste + Production
        if col == "Waste %":
            continue
        fit = _linfit_y_on_prod(df_all, col)
        if fit is None:
            continue
        a, b, resid_std = fit
        z = 1.28  # ~80% band
        y_lo = a * prod_lo + b
        y_hi = a * prod_hi + b
        lo = min(y_lo, y_hi) - z * resid_std
        hi = max(y_lo, y_hi) + z * resid_std
        midv = a * mid_prod + b
        if col.endswith("%"):
            lo, midv, hi = _clip_pct(lo), _clip_pct(midv), _clip_pct(hi)
        else:
            lo, midv, hi = max(0.0, lo), max(0.0, midv), max(0.0, hi)
        out[col] = {"mid": float(midv), "lo": float(lo), "hi": float(hi)}

    # Derive Waste % from predicted Waste & Production, if both available
    if out.get("Waste (Tons)"):
        w = out["Waste (Tons)"]
        if prod_lo > 0 and prod_hi > 0:
            lo_share = (w["lo"] / (prod_hi + w["lo"])) * 100.0 if (prod_hi + w["lo"]) > 0 else 0.0
            hi_share = (w["hi"] / (prod_lo + w["hi"])) * 100.0 if (prod_lo + w["hi"]) > 0 else 0.0
            mid_share = (w["mid"] / (mid_prod + w["mid"])) * 100.0 if (mid_prod + w["mid"]) > 0 else 0.0
            out["Waste %"] = {
                "mid": _clip_pct(mid_share),
                "lo": _clip_pct(min(lo_share, hi_share)),
                "hi": _clip_pct(max(lo_share, hi_share)),
            }

    # Include Production identity prediction
    out["Production (Tons)"] = {"mid": mid_prod, "lo": float(prod_lo), "hi": float(prod_hi)}
    return out

def section_forecast(df_all: pd.DataFrame):
    """
    Forecast page:
    - Input a Production (Tons) MIN–MAX
    - Predict ranges for key KPIs using linear fit vs Production + residual spread
    - Show cards + charts with forecast markers
    - Provide download row via _download_row
    """
    if df_all.empty:
        st.info("Upload KPI file to view Forecast.")
        return {}

    charts: Dict[str, go.Figure] = {}

    # --- Next-month label
    if "Month" in df_all.columns and df_all["Month"].notna().any():
        last_dt = pd.to_datetime(df_all["Month"].max())
        next_dt = (last_dt + pd.offsets.MonthBegin(1))
        next_month_label = next_dt.strftime("%b %Y")
        next_month_short = next_dt.strftime("%b")
    else:
        next_dt = None
        next_month_label = "Next Month"
        next_month_short = "Next"

    # --- Default production range (last value ±10%)
    hist_prod = df_all["Production (Tons)"].dropna().astype(float) if "Production (Tons)" in df_all.columns else pd.Series(dtype=float)
    last_prod = float(hist_prod.iloc[-1]) if not hist_prod.empty else 0.0
    default_lo = round(last_prod * 0.90, 2) if last_prod > 0 else 0.0
    default_hi = round(last_prod * 1.10, 2) if last_prod > 0 else 0.0

    # --- Controls
    st.markdown(f'<div class="ppk-panel"><div class="ppk-titlebar">Forecast setup</div><div class="ppk-accent"></div>', unsafe_allow_html=True)
    c1, c2 = st.columns(2)
    with c1:
        prod_lo = st.number_input("Min Production (Tons)", min_value=0.0, value=float(default_lo), step=10.0, key="forecast_prod_min_new")
    with c2:
        prod_hi = st.number_input("Max Production (Tons)", min_value=0.0, value=float(default_hi if default_hi >= default_lo else default_lo), step=10.0, key="forecast_prod_max_new")
    st.markdown("</div>", unsafe_allow_html=True)

    if prod_hi < prod_lo:
        st.warning("Max Production must be ≥ Min Production.")
        return charts
    if prod_hi <= 0:
        st.info("Enter a positive Production range to generate the forecast.")
        return charts

    # --- Predictions
    preds = _predict_range_from_prod(df_all, prod_lo, prod_hi)

    # --- Cards (mid + range)
    st.markdown(f"#### Forecasted KPIs — {next_month_label}")
    tile_cols = st.columns(6, gap="small")
    order_cards = [
        ("Production (Tons)", False),
        ("Number of Accidents", False),
        ("Waste %", True),
        ("Total Down Time (Hours)", False),
        ("OEE %", True),
        ("Reclaim (Tons)", False),
    ]
    for i, (kpi, is_pct) in enumerate(order_cards):
        p = preds.get(kpi, {})
        mid, lo, hi = p.get("mid"), p.get("lo"), p.get("hi")
        band = f"{fmt_num(lo, is_pct)} – {fmt_num(hi, is_pct)}" if (lo is not None and hi is not None) else "—"
        with tile_cols[i]:
            st.markdown(
                f"""
                <div class="metric-card">
                  <div class="metric-caption">{kpi} (forecast)</div>
                  <div class="big-value">{fmt_num(mid, is_pct)}</div>
                  <span class="delta-badge">Range: {band}</span>
                </div>
                """,
                unsafe_allow_html=True
            )

    # --- Build plotting context; append forecast midpoint row
    ctx = df_year_default.copy() if 'df_year_default' in globals() else pd.DataFrame()
    if not ctx.empty and next_dt is not None:
        add_row = {"Month": next_dt, "Year": int(next_dt.year), "MonthNum": int(next_dt.month), "MonthName": next_month_short}
        for k, v in preds.items():
            if v.get("mid") is not None:
                add_row[k] = v["mid"]
        ctx = pd.concat([ctx, pd.DataFrame([add_row])], ignore_index=True).sort_values("MonthNum")
        x_forecast_label = next_month_short
    else:
        x_forecast_label = ctx.iloc[-1]["MonthName"] if not ctx.empty else next_month_short

    def _maybe_mark(fig: go.Figure, key: str):
        v = preds.get(key, {})
        if v and v.get("mid") is not None:
            add_range_marker(fig, x_forecast_label, v["lo"], v["mid"], v["hi"], name=f"{key} (range)")

    figs: List[go.Figure] = []

    if "Production (Tons)" in ctx.columns:
        st.markdown(f'<div class="ppk-panel"><div class="ppk-titlebar">Production — Monthly</div><div class="ppk-accent"></div>', unsafe_allow_html=True)
        f = line_fig(ctx, "Production (Tons)", "", COLORWAY[0], pct=False)
        _maybe_mark(f, "Production (Tons)")
        charts["forecast_production"] = f; figs.append(f)
        st.plotly_chart(f, use_container_width=True); st.markdown("</div>", unsafe_allow_html=True)

    if "OEE %" in ctx.columns:
        st.markdown(f'<div class="ppk-panel"><div class="ppk-titlebar">OEE% — Monthly</div><div class="ppk-accent"></div>', unsafe_allow_html=True)
        f = area_fig(ctx, "OEE %", "", COLORWAY[3], pct=True)
        _maybe_mark(f, "OEE %")
        charts["forecast_oee"] = f; figs.append(f)
        st.plotly_chart(f, use_container_width=True); st.markdown("</div>", unsafe_allow_html=True)

    if {"Production (Tons)", "Waste (Tons)"}.issubset(ctx.columns):
        st.markdown(f'<div class="ppk-panel"><div class="ppk-titlebar">Total Production by Month (Stacked) — with Forecast</div><div class="ppk-accent"></div>', unsafe_allow_html=True)
        f = stacked_bar_fig(ctx, "Production (Tons)", "Waste (Tons)", "", "#60A5FA", "#F59E0B")
        charts["forecast_stacked"] = f; figs.append(f)
        st.plotly_chart(f, use_container_width=True); st.markdown("</div>", unsafe_allow_html=True)

    if "Total Down Time (Hours)" in ctx.columns:
        st.markdown(f'<div class="ppk-panel"><div class="ppk-titlebar">Downtime — Monthly</div><div class="ppk-accent"></div>', unsafe_allow_html=True)
        f = area_fig(ctx, "Total Down Time (Hours)", "", "#06B6D4")
        _maybe_mark(f, "Total Down Time (Hours)")
        charts["forecast_downtime"] = f; figs.append(f)
        st.plotly_chart(f, use_container_width=True); st.markdown("</div>", unsafe_allow_html=True)

    if "Reclaim (Tons)" in ctx.columns:
        st.markdown(f'<div class="ppk-panel"><div class="ppk-titlebar">Reclaim — Monthly</div><div class="ppk-accent"></div>', unsafe_allow_html=True)
        f = lollipop_fig(ctx, "Reclaim (Tons)", "", TEAL)
        _maybe_mark(f, "Reclaim (Tons)")
        charts["forecast_reclaim"] = f; figs.append(f)
        st.plotly_chart(f, use_container_width=True); st.markdown("</div>", unsafe_allow_html=True)

    if "Number of Accidents" in ctx.columns:
        st.markdown(f'<div class="ppk-panel"><div class="ppk-titlebar">Number of Accidents — Monthly</div><div class="ppk-accent"></div>', unsafe_allow_html=True)
        f = bar_fig(ctx, "Number of Accidents", "", COLORWAY[5])
        _maybe_mark(f, "Number of Accidents")
        charts["forecast_accidents"] = f; figs.append(f)
        st.plotly_chart(f, use_container_width=True); st.markdown("</div>", unsafe_allow_html=True)

    if "Waste %" in ctx.columns:
        st.markdown(f'<div class="ppk-panel"><div class="ppk-titlebar">Waste % — Monthly</div><div class="ppk-accent"></div>', unsafe_allow_html=True)
        f = line_fig(ctx, "Waste %", "", COLORWAY[2], pct=True)
        _maybe_mark(f, "Waste %")
        charts["forecast_wastepct"] = f; figs.append(f)
        st.plotly_chart(f, use_container_width=True); st.markdown("</div>", unsafe_allow_html=True)

    if preds.get("Waste (Tons)") and preds.get("Production (Tons)"):
        donut_df = pd.DataFrame({
            "Category": ["Good Production", "Waste"],
            "Tons": [preds["Production (Tons)"]["mid"], preds["Waste (Tons)"]["mid"]],
        })
        st.markdown(f'<div class="ppk-panel"><div class="ppk-titlebar">Predicted Composition — {next_month_label} (midpoint)</div><div class="ppk-accent"></div>', unsafe_allow_html=True)
        f = px.pie(
            donut_df, names="Category", values="Tons", hole=0.55,
            color="Category",
            color_discrete_map={"Good Production": "#34D399", "Waste": "#EF4444"}
        )
        f.update_traces(textinfo="percent+label", pull=[0, 0.06])
        charts["donut_f"] = _apply_theme(f); figs.append(charts["donut_f"])
        st.plotly_chart(charts["donut_f"], use_container_width=True)
        st.markdown("</div>", unsafe_allow_html=True)

    forecast_rows = []
    for k in ["Production (Tons)", "Waste (Tons)", "Waste %", "Total Down Time (Hours)", "OEE %", "Reclaim (Tons)", "Number of Accidents"]:


        v = preds.get(k, {})
        if v:
            forecast_rows.append({
                "KPI": k,
                "Production_Min": prod_lo,
                "Production_Max": prod_hi,
                "Pred_Lo": v.get("lo"),
                "Pred_Mid": v.get("mid"),
                "Pred_Hi": v.get("hi"),
                "IsPercent": k.endswith("%"),
                "Forecast_Month": next_month_label,
                "Scope": "All years",
            })
    forecast_table = pd.DataFrame(forecast_rows)

    narrative = [
        f"Forecast month: {next_month_label}.",
        "Fit scope: All available history (linear vs Production).",
        f"Production input range: {fmt_num(prod_lo)} – {fmt_num(prod_hi)}.",
        "Bands combine production-range effect + ~80% residual spread.",
    ]
    _download_row("Forecast", forecast_table, figs, narrative)
    return charts

    if df_all is None or df_all.empty or "Year" not in df_all.columns:
        return None, None, df_all

    years = sorted(df_all["Year"].dropna().astype(int).unique().tolist())
    default_year = years[-1] if years else None

    c1, c2 = st.columns([1, 1])
    with c1:
        sel_year = st.selectbox(
            f"{label_prefix}Year",
            options=years,
            index=(years.index(default_year) if default_year in years else 0),
            key=f"ym_year_{label_prefix.replace(' ', '_').lower()}"
        )

    df_year = df_all[df_all["Year"] == sel_year].copy().sort_values("MonthNum")
    months_df = df_year[["MonthNum", "MonthName"]].dropna().drop_duplicates().sort_values("MonthNum")
    month_name_by_num = dict(zip(months_df["MonthNum"].astype(int), months_df["MonthName"]))

    with c2:
        month_options = ["All months"] + [f"{m:02d} — {month_name_by_num.get(m, '')}" for m in months_df["MonthNum"].astype(int)]
        sel_month_label = st.selectbox(
            f"{label_prefix}Month",
            options=month_options,
            index=0,
            key=f"ym_month_{label_prefix.replace(' ', '_').lower()}"
        )

    if sel_month_label == "All months":
        return sel_year, None, df_year

    sel_monthnum = int(sel_month_label.split("—")[0].strip())
    df_trim = df_year[df_year["MonthNum"] <= sel_monthnum].copy()
    return sel_year, sel_monthnum, df_trim

# =========================
# PAGE ROUTER
# =========================
# =========================
# PAGE ROUTER
# =========================
active = st.session_state.get("active_page", "Home")

if active == "Home":
    section_home(df, df_year_default)

elif active == "Monthly Update":
    section_monthly_update(df)

elif active == "Forecast":
    section_forecast(df)

elif active == "Waste Analytics":
    # page-level Year/Month for KPI dataframe
    _, _, df_for_tab = pick_year_month(df, "Waste Analytics — ")
    if df_for_tab is None or df_for_tab.empty:
        st.info("Upload KPI data to view Waste Analytics.")
    else:
        section_waste(df_for_tab)

elif active == "Breakup":
    section_breakup_combined()

elif active == "KPIs Comparison":
    # compare using the KPI dataframe filtered to YTD up to the picked month
    _, _, df_for_tab = pick_year_month(df, "KPIs Comparison — ")
    if df_for_tab is None or df_for_tab.empty:
        st.info("Upload KPI data to view KPIs Comparison.")
    else:
        section_kpi_comparison(df_for_tab)

elif active == "Reports":
    # independent pickers for each dataset used in Reports
    _, _, df_kpi_for_tab = pick_year_month(df,   "Reports (KPI) — ")
    _, _, df_wb_for_tab  = pick_year_month(df_wb, "Reports (Waste) — ")
    _, _, df_dt_for_tab  = pick_year_month(df_dt, "Reports (Downtime) — ")
    section_reports(df_kpi_for_tab, df_wb_for_tab, df_dt_for_tab)

# =========================
# FOOTER BAR (fixed) + DEV MODAL
# =========================
def _get_qp(key: str) -> str:
    try:
        qp = st.query_params
        return qp.get(key, [""])[0] if isinstance(qp.get(key), list) else (qp.get(key) or "")
    except Exception:
        qp = st.experimental_get_query_params()
        return (qp.get(key, [""]) or [""])[0]

def _strip_qp() -> str:
    # go back to same path with no query so modal closes
    return "?"


def render_dev_modal():
    if not st.session_state.get("show_dev_modal", False):
        return

    st.markdown(
        f"""
<div class="modal-overlay">
  <div class="modal-card">
    <div class="modal-title">About the Developer</div>
    <div class="modal-muted">Built by <b>{DEV_NAME}</b> — {DEV_ROLE}</div>

    <p><b>SARKS — Secure Access Role Knowledge System</b></p>

    <p>This model was developed by <b>Sheikh Abdul Rehman</b> during his internship at <b>SARKS</b>. 
       It follows SARKS’s visual identity (logo, colors and typography) and presents factory KPIs 
       in a clean, decision-ready format.</p>

    <p><b>What this model does (plain language):</b></p>
    <ul>
      <li><b>Home:</b> Key yearly totals/averages with easy month-to-month comparison.</li>
      <li><b>Breakup:</b> Shows <i>where</i> waste and downtime come from (factors, Pareto, shares).</li>
      <li><b>Forecast:</b> Gives a <i>simple</i> next-month estimate with an uncertainty band.</li>
      <li><b>Reports:</b> Clean tables and overview charts for printing/sharing.</li>
    </ul>

    <p class="modal-muted">
      Tech: Streamlit • Plotly • Pandas/NumPy • (optional) Kaleido for image export • ReportLab (PDF) • python-pptx (PPT).
    </p>

    <p class="modal-muted">
      Connect on LinkedIn: 
      <a href="{DEV_LINK}" target="_blank" rel="noopener">LinkedIn • Sheikh Abdul Rehman bin Khalid Sharif</a>
    </p>

    <div class="modal-actions">
      <!-- Close returns to the same path without ?dev=1 -->
      <a href="{_strip_qp()}"><button class="modal-close-btn">Close</button></a>
    </div>
  </div>
</div>
        """,
        unsafe_allow_html=True,
    )
# --- About page router (fixed)
dev_qs = st.query_params.get("dev")
if dev_qs == "1":
    st.query_params.clear()
    render_about_developer_page()
    st.stop()

# Always show footer bar (only)
render_footer_bar()


